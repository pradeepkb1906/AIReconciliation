"""
title: AI Reconciliation
author: IBM Consulting Advantage
version: 1.0.0
release: V1.0 (2026-04-25T05:07:49Z) — first end-to-end working release.
description: Stepwise in-iframe reconciliation wizard. Accepts CSV/TSV/XLSX/XLS/DOCX/PPTX/PDF/JSON uploads, runs deterministic matching in the browser, and generates XLSX/CSV/PDF/DOCX outputs accepted by USA and EU regulators (SOX, BCBS 239, EMIR, MiFID II, Solvency II, HIPAA, 21 CFR Part 11, GDPR, DORA). Falls back to server-side ooXML generation when CDN libraries are blocked. Requires iframe Sandbox Allow Same Origin in Open WebUI Settings -> Interface.
"""

from __future__ import annotations

import base64
import html as htmlmod
import io
import json
import logging
import re
import time
import traceback
import uuid

log = logging.getLogger("reconciliation_tool")
log.setLevel(logging.INFO)
from typing import Any, Dict, List, Literal, Optional

from fastapi.responses import HTMLResponse
from pydantic import BaseModel, Field

_BUILD = "1.0.0"
_RELEASE_TS = "2026-04-25T05:07:49Z"  # V1.0 — first working end-to-end release

# ---------------------------------------------------------------------------
# IBM Light Navy Blue theme
# ---------------------------------------------------------------------------

IBM_NAVY = "#002D74"
IBM_NAVY_DARK = "#001F52"
IBM_BLUE = "#0043CE"
IBM_BLUE_LIGHT = "#4589FF"
IBM_SURFACE = "#EDF5FF"
IBM_SURFACE_2 = "#D0E2FF"
IBM_TEXT = "#161616"
IBM_TEXT_2 = "#525252"
IBM_BORDER = "#C6C6C6"
IBM_SUCCESS = "#24A148"
IBM_WARN = "#F1C21B"
IBM_DANGER = "#DA1E28"

THEME_CSS = f"""
:root {{
  --ibm-navy: {IBM_NAVY};
  --ibm-navy-dark: {IBM_NAVY_DARK};
  --ibm-blue: {IBM_BLUE};
  --ibm-blue-light: {IBM_BLUE_LIGHT};
  --ibm-surface: {IBM_SURFACE};
  --ibm-surface-2: {IBM_SURFACE_2};
  --ibm-text: {IBM_TEXT};
  --ibm-text-2: {IBM_TEXT_2};
  --ibm-border: {IBM_BORDER};
  --ibm-success: {IBM_SUCCESS};
  --ibm-warn: {IBM_WARN};
  --ibm-danger: {IBM_DANGER};
  --font-sans: 'IBM Plex Sans', system-ui, -apple-system, 'Segoe UI', Roboto, sans-serif;
  --font-mono: 'IBM Plex Mono', 'SF Mono', Menlo, Consolas, monospace;
  --radius: 4px;
  --radius-lg: 8px;
  color-scheme: light;
}}
* {{ box-sizing: border-box; }}
html, body {{ margin: 0; padding: 0; background: #FFFFFF; color: var(--ibm-text); font-family: var(--font-sans); font-size: 14px; line-height: 1.5; }}
a {{ color: var(--ibm-blue); }}
.ibm-shell {{ padding: 20px 24px 32px; max-width: 100%; }}
.ibm-masthead {{ display: flex; align-items: center; justify-content: space-between; gap: 16px; padding: 14px 20px; background: var(--ibm-navy); color: #FFFFFF; border-radius: var(--radius) var(--radius) 0 0; }}
.ibm-masthead .brand {{ display: flex; align-items: center; gap: 12px; font-weight: 600; letter-spacing: 0.01em; }}
.ibm-masthead .brand .dot {{ width: 10px; height: 10px; background: var(--ibm-blue-light); border-radius: 50%; display: inline-block; }}
.ibm-masthead .meta {{ font-size: 12px; opacity: 0.85; font-family: var(--font-mono); }}
.ibm-subbar {{ display: flex; flex-wrap: wrap; gap: 10px; align-items: center; padding: 10px 20px; background: var(--ibm-surface); border-left: 4px solid var(--ibm-blue); color: var(--ibm-text); font-size: 13px; }}
.ibm-chip {{ display: inline-flex; align-items: center; gap: 6px; padding: 4px 10px; background: #FFFFFF; border: 1px solid var(--ibm-border); border-radius: 999px; font-size: 12px; }}
.ibm-chip.primary {{ background: var(--ibm-navy); color: #FFFFFF; border-color: var(--ibm-navy); }}
.ibm-chip.warn {{ background: #FFF8E1; border-color: var(--ibm-warn); }}
.ibm-chip.danger {{ background: #FFF1F1; border-color: var(--ibm-danger); color: var(--ibm-danger); }}
.ibm-chip.success {{ background: #DEFBE6; border-color: var(--ibm-success); color: #0E6027; }}
.ibm-card {{ background: #FFFFFF; border: 1px solid var(--ibm-border); border-radius: var(--radius); margin-top: 16px; overflow: hidden; }}
.ibm-card > header {{ padding: 12px 16px; background: #F4F4F4; border-bottom: 1px solid var(--ibm-border); font-weight: 600; color: var(--ibm-navy); display: flex; align-items: center; justify-content: space-between; gap: 10px; }}
.ibm-card > .body {{ padding: 14px 16px; }}
.ibm-kpis {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(160px, 1fr)); gap: 12px; }}
.ibm-kpi {{ background: var(--ibm-surface); border: 1px solid var(--ibm-surface-2); border-radius: var(--radius); padding: 12px 14px; }}
.ibm-kpi .v {{ font-size: 26px; font-weight: 600; color: var(--ibm-navy); letter-spacing: -0.01em; }}
.ibm-kpi .l {{ font-size: 12px; color: var(--ibm-text-2); margin-top: 2px; }}
.ibm-kpi.ok .v {{ color: var(--ibm-success); }}
.ibm-kpi.warn .v {{ color: #B28600; }}
.ibm-kpi.danger .v {{ color: var(--ibm-danger); }}
.ibm-table {{ width: 100%; border-collapse: collapse; font-size: 13px; }}
.ibm-table th {{ background: var(--ibm-navy); color: #FFFFFF; font-weight: 500; text-align: left; padding: 8px 12px; border-bottom: 2px solid var(--ibm-navy-dark); }}
.ibm-table td {{ padding: 8px 12px; border-bottom: 1px solid #E0E0E0; vertical-align: top; }}
.ibm-table tr:nth-child(even) td {{ background: #FAFAFA; }}
.ibm-table tr:hover td {{ background: var(--ibm-surface); }}
.ibm-table td.num, .ibm-table th.num {{ text-align: right; font-variant-numeric: tabular-nums; font-family: var(--font-mono); }}
.ibm-table td.delta-pos {{ color: var(--ibm-success); }}
.ibm-table td.delta-neg {{ color: var(--ibm-danger); }}
.ibm-table td.key {{ font-family: var(--font-mono); font-size: 12px; color: var(--ibm-navy); }}
.ibm-tab-scroll {{ overflow-x: auto; }}
.ibm-btn {{ display: inline-flex; align-items: center; gap: 6px; padding: 8px 14px; background: var(--ibm-navy); color: #FFFFFF; border: none; border-radius: var(--radius); font: 500 13px var(--font-sans); cursor: pointer; transition: background 120ms; }}
.ibm-btn:hover {{ background: var(--ibm-navy-dark); }}
.ibm-btn.secondary {{ background: #FFFFFF; color: var(--ibm-navy); border: 1px solid var(--ibm-navy); }}
.ibm-btn.secondary:hover {{ background: var(--ibm-surface); }}
.ibm-btn[disabled] {{ opacity: 0.5; cursor: not-allowed; }}
.ibm-dl-bar {{ position: sticky; top: 0; z-index: 10; display: flex; gap: 8px; padding: 10px 20px; background: rgba(255,255,255,0.96); backdrop-filter: blur(4px); border-bottom: 1px solid var(--ibm-border); align-items: center; flex-wrap: wrap; }}
.ibm-dl-bar .label {{ font-weight: 600; color: var(--ibm-navy); margin-right: 6px; }}
.ibm-dl-bar .status {{ margin-left: auto; font-size: 12px; color: var(--ibm-text-2); }}
.ibm-footer {{ padding: 12px 20px; margin-top: 20px; border-top: 2px solid var(--ibm-navy); font-size: 11px; color: var(--ibm-text-2); display: flex; justify-content: space-between; gap: 12px; flex-wrap: wrap; }}
.ibm-loading {{ display: flex; flex-direction: column; align-items: center; justify-content: center; padding: 40px 20px; color: var(--ibm-text-2); }}
.ibm-loading .dots span {{ display: inline-block; width: 10px; height: 10px; margin: 0 3px; border-radius: 50%; background: var(--ibm-blue); animation: ibm-bounce 1.2s ease-in-out infinite; }}
.ibm-loading .dots span:nth-child(2) {{ animation-delay: 0.2s; }}
.ibm-loading .dots span:nth-child(3) {{ animation-delay: 0.4s; }}
@keyframes ibm-bounce {{ 0%, 80%, 100% {{ transform: scale(0.4); opacity: 0.4; }} 40% {{ transform: scale(1); opacity: 1; }} }}
.ibm-section-title {{ font-size: 16px; font-weight: 600; color: var(--ibm-navy); margin: 22px 0 8px; display: flex; align-items: center; gap: 8px; }}
.ibm-section-title .bar {{ width: 4px; height: 18px; background: var(--ibm-blue); border-radius: 2px; }}
.ibm-narrative {{ background: #FAFAFA; border-left: 3px solid var(--ibm-blue); padding: 12px 16px; color: var(--ibm-text); font-size: 13.5px; border-radius: 0 var(--radius) var(--radius) 0; }}
.ibm-narrative p {{ margin: 6px 0; }}
.ibm-regtag {{ display: inline-block; padding: 2px 8px; background: var(--ibm-surface-2); color: var(--ibm-navy-dark); border-radius: 999px; font-size: 11px; font-weight: 500; margin-right: 4px; margin-bottom: 2px; }}
.ibm-toast {{ position: fixed; right: 18px; bottom: 18px; padding: 10px 16px; background: var(--ibm-navy); color: #FFFFFF; border-radius: var(--radius); font-size: 13px; box-shadow: 0 4px 14px rgba(0,0,0,0.18); opacity: 0; transform: translateY(10px); transition: opacity 160ms, transform 160ms; pointer-events: none; z-index: 50; }}
.ibm-toast.show {{ opacity: 1; transform: translateY(0); }}
.ibm-toast.err {{ background: var(--ibm-danger); }}
@media print {{
  .ibm-dl-bar, .ibm-btn {{ display: none !important; }}
  .ibm-masthead, .ibm-card > header {{ -webkit-print-color-adjust: exact; print-color-adjust: exact; }}
}}
"""

# ---------------------------------------------------------------------------
# Pastel wizard theme (IBM Light Navy, very light and pleasant)
# ---------------------------------------------------------------------------

WIZARD_CSS = """
:root {
  --wz-bg: #F4F8FC;
  --wz-surface: #FFFFFF;
  --wz-surface-soft: #EAF2FB;
  --wz-surface-hover: #DCE8F6;
  --wz-pastel: #E3EEFA;
  --wz-pastel-2: #D1E2F3;
  --wz-ink: #0F2A44;
  --wz-ink-soft: #3B5472;
  --wz-muted: #6B7F97;
  --wz-accent: #4178BE;
  --wz-accent-strong: #2A5C9A;
  --wz-accent-soft: #BBD4EE;
  --wz-border: #C7D9EC;
  --wz-ring: rgba(65, 120, 190, 0.28);
  --wz-ok: #2E7D32;
  --wz-ok-soft: #E1F1E4;
  --wz-warn: #B25F00;
  --wz-warn-soft: #FFF4E0;
  --wz-err: #B0241D;
  --wz-err-soft: #FBE6E4;
}
.wz-root { background: var(--wz-bg); min-height: 100vh; padding: 24px 20px 60px; font-family: 'IBM Plex Sans', system-ui, -apple-system, 'Segoe UI', Roboto, sans-serif; color: var(--wz-ink); }
.wz-shell { max-width: 1160px; margin: 0 auto; }
.wz-head { display: flex; align-items: center; justify-content: space-between; gap: 16px; padding: 18px 22px; background: linear-gradient(135deg, #F7FAFE 0%, var(--wz-surface-soft) 100%); border: 1px solid var(--wz-border); border-radius: 12px; box-shadow: 0 1px 2px rgba(15,42,68,0.04); }
.wz-head h1 { margin: 0; font-size: 18px; font-weight: 600; color: var(--wz-ink); letter-spacing: -0.01em; display: flex; align-items: center; gap: 12px; }
.wz-head h1 .wz-dot { width: 10px; height: 10px; background: var(--wz-accent); border-radius: 50%; display: inline-block; box-shadow: 0 0 0 4px var(--wz-pastel); }
.wz-head .wz-meta { font-size: 12px; color: var(--wz-muted); font-family: 'IBM Plex Mono', 'SF Mono', Menlo, monospace; }

.wz-stepper { display: flex; align-items: stretch; gap: 0; margin: 18px 0 8px; padding: 14px 16px; background: var(--wz-surface); border: 1px solid var(--wz-border); border-radius: 12px; overflow-x: auto; }
.wz-step { display: flex; align-items: center; gap: 10px; flex: 1 1 0; min-width: 180px; padding: 6px 12px; cursor: pointer; border-radius: 8px; color: var(--wz-ink-soft); transition: background 140ms ease; }
.wz-step:hover { background: var(--wz-pastel); }
.wz-step .wz-bullet { width: 28px; height: 28px; border-radius: 50%; display: grid; place-items: center; background: var(--wz-pastel-2); color: var(--wz-accent-strong); font-size: 13px; font-weight: 600; border: 1px solid var(--wz-accent-soft); flex-shrink: 0; }
.wz-step.active .wz-bullet { background: var(--wz-accent); color: #FFFFFF; border-color: var(--wz-accent); box-shadow: 0 0 0 4px var(--wz-ring); }
.wz-step.done .wz-bullet { background: var(--wz-ok); color: #FFFFFF; border-color: var(--wz-ok); }
.wz-step.done .wz-bullet::before { content: ''; width: 9px; height: 5px; border-left: 2px solid #fff; border-bottom: 2px solid #fff; transform: rotate(-45deg) translate(0, -1px); position: absolute; }
.wz-step.done .wz-bullet { position: relative; color: transparent; }
.wz-step .wz-txt { display: flex; flex-direction: column; line-height: 1.2; }
.wz-step .wz-tag { font-size: 11px; color: var(--wz-muted); letter-spacing: 0.02em; text-transform: uppercase; font-weight: 500; }
.wz-step .wz-lbl { font-size: 13.5px; color: var(--wz-ink); font-weight: 500; }
.wz-step.active .wz-lbl { color: var(--wz-accent-strong); font-weight: 600; }
.wz-step .wz-bar { display: none; }

.wz-panel { background: var(--wz-surface); border: 1px solid var(--wz-border); border-radius: 12px; padding: 22px 26px; margin-top: 16px; box-shadow: 0 1px 2px rgba(15,42,68,0.04); }
.wz-panel h2 { margin: 0 0 4px; font-size: 17px; font-weight: 600; color: var(--wz-ink); }
.wz-panel .wz-sub { margin: 0 0 20px; color: var(--wz-muted); font-size: 13.5px; }
.wz-grid { display: grid; grid-template-columns: repeat(auto-fit, minmax(260px, 1fr)); gap: 14px; }
.wz-field { display: flex; flex-direction: column; gap: 6px; }
.wz-field label { font-size: 12.5px; font-weight: 500; color: var(--wz-ink-soft); letter-spacing: 0.01em; }
.wz-field .wz-hint { font-size: 11.5px; color: var(--wz-muted); margin-top: 2px; }
.wz-input, .wz-select, .wz-textarea { font: inherit; padding: 10px 12px; border: 1px solid var(--wz-border); background: #FBFDFF; border-radius: 8px; color: var(--wz-ink); outline: none; transition: border-color 140ms, box-shadow 140ms, background 140ms; }
.wz-input:focus, .wz-select:focus, .wz-textarea:focus { border-color: var(--wz-accent); box-shadow: 0 0 0 3px var(--wz-ring); background: #FFFFFF; }
.wz-select { appearance: none; -webkit-appearance: none; background-image: linear-gradient(45deg, transparent 50%, var(--wz-accent) 50%), linear-gradient(135deg, var(--wz-accent) 50%, transparent 50%); background-position: calc(100% - 18px) 50%, calc(100% - 12px) 50%; background-size: 6px 6px; background-repeat: no-repeat; padding-right: 32px; }

.wz-drop { border: 2px dashed var(--wz-accent-soft); border-radius: 12px; padding: 26px 20px; background: linear-gradient(180deg, #F7FBFF 0%, #EEF5FC 100%); text-align: center; transition: border-color 140ms, background 140ms, transform 140ms; cursor: pointer; }
.wz-drop:hover, .wz-drop.wz-hover { border-color: var(--wz-accent); background: #E8F1FA; }
.wz-drop.wz-has-files { border-style: solid; border-color: var(--wz-accent); background: #F1F7FD; }
.wz-drop .wz-icon { width: 44px; height: 44px; border-radius: 50%; background: var(--wz-pastel-2); display: grid; place-items: center; margin: 0 auto 10px; color: var(--wz-accent-strong); font-size: 22px; font-weight: 700; }
.wz-drop .wz-title { font-weight: 600; color: var(--wz-ink); font-size: 14.5px; margin-bottom: 4px; }
.wz-drop .wz-sub { color: var(--wz-muted); font-size: 12.5px; margin-bottom: 10px; }
.wz-drop .wz-browse { display: inline-block; padding: 7px 14px; background: var(--wz-accent); color: #FFFFFF; border-radius: 999px; font-size: 12.5px; font-weight: 500; letter-spacing: 0.01em; }
.wz-drop input[type=file] { display: none; }

.wz-filelist { margin-top: 12px; display: flex; flex-direction: column; gap: 6px; }
.wz-file { display: flex; align-items: center; gap: 10px; padding: 8px 12px; background: var(--wz-surface-soft); border: 1px solid var(--wz-border); border-radius: 8px; font-size: 12.5px; color: var(--wz-ink); }
.wz-file .wz-ext { padding: 2px 8px; background: var(--wz-accent); color: #FFFFFF; border-radius: 4px; font-size: 10.5px; font-weight: 600; letter-spacing: 0.02em; text-transform: uppercase; }
.wz-file .wz-name { flex: 1; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }
.wz-file .wz-status { font-size: 11.5px; color: var(--wz-muted); }
.wz-file .wz-status.ok { color: var(--wz-ok); }
.wz-file .wz-status.err { color: var(--wz-err); }
.wz-file .wz-rem { background: transparent; border: none; color: var(--wz-muted); cursor: pointer; font-size: 14px; line-height: 1; padding: 2px 6px; border-radius: 4px; }
.wz-file .wz-rem:hover { background: var(--wz-err-soft); color: var(--wz-err); }

.wz-progress { height: 8px; background: var(--wz-pastel); border-radius: 999px; overflow: hidden; margin-top: 8px; }
.wz-progress-bar { height: 100%; background: linear-gradient(90deg, var(--wz-accent) 0%, var(--wz-accent-strong) 100%); border-radius: 999px; transition: width 180ms ease; width: 0%; }
.wz-progress-row { display: flex; align-items: center; gap: 12px; margin-top: 10px; font-size: 12.5px; color: var(--wz-ink-soft); }
.wz-progress-row .wz-lbl { min-width: 150px; }
.wz-progress-row .wz-val { margin-left: auto; color: var(--wz-muted); font-family: 'IBM Plex Mono', monospace; font-size: 12px; }

.wz-preview { margin-top: 14px; border: 1px solid var(--wz-border); border-radius: 10px; overflow: hidden; }
.wz-preview header { background: var(--wz-surface-soft); padding: 10px 14px; font-size: 12.5px; color: var(--wz-ink-soft); display: flex; justify-content: space-between; align-items: center; gap: 10px; flex-wrap: wrap; }
.wz-preview header .wz-tabselect { display: flex; gap: 6px; flex-wrap: wrap; }
.wz-preview header .wz-tab { padding: 4px 10px; background: #FFFFFF; border: 1px solid var(--wz-border); border-radius: 999px; font-size: 11.5px; cursor: pointer; color: var(--wz-ink-soft); }
.wz-preview header .wz-tab.active { background: var(--wz-accent); color: #FFFFFF; border-color: var(--wz-accent); }
.wz-preview .wz-scroll { max-height: 260px; overflow: auto; background: #FFFFFF; }
.wz-preview table { width: 100%; border-collapse: collapse; font-size: 12.5px; }
.wz-preview th { background: #F5F9FD; color: var(--wz-ink-soft); font-weight: 500; text-align: left; padding: 8px 10px; border-bottom: 1px solid var(--wz-border); position: sticky; top: 0; z-index: 2; }
.wz-preview td { padding: 7px 10px; border-bottom: 1px solid #EEF3F8; color: var(--wz-ink); }
.wz-preview tr:nth-child(even) td { background: #FBFDFF; }

.wz-actions { display: flex; justify-content: space-between; align-items: center; gap: 12px; margin-top: 22px; padding-top: 18px; border-top: 1px solid var(--wz-border); }
.wz-actions .wz-hint { font-size: 12px; color: var(--wz-muted); }
.wz-btn { display: inline-flex; align-items: center; gap: 8px; padding: 10px 18px; border: 1px solid var(--wz-accent); background: var(--wz-accent); color: #FFFFFF; border-radius: 8px; font: 500 13.5px 'IBM Plex Sans', sans-serif; cursor: pointer; transition: background 140ms, transform 80ms, box-shadow 140ms; }
.wz-btn:hover { background: var(--wz-accent-strong); border-color: var(--wz-accent-strong); }
.wz-btn:active { transform: translateY(1px); }
.wz-btn:disabled { opacity: 0.5; cursor: not-allowed; }
.wz-btn.wz-secondary { background: #FFFFFF; color: var(--wz-accent-strong); border-color: var(--wz-accent-soft); }
.wz-btn.wz-secondary:hover { background: var(--wz-pastel); border-color: var(--wz-accent); }
.wz-btn.wz-ghost { background: transparent; color: var(--wz-ink-soft); border-color: transparent; }
.wz-btn.wz-ghost:hover { background: var(--wz-pastel); color: var(--wz-ink); }

.wz-kpis { display: grid; grid-template-columns: repeat(auto-fit, minmax(160px, 1fr)); gap: 12px; margin-top: 8px; }
.wz-kpi { background: linear-gradient(160deg, #FFFFFF 0%, var(--wz-surface-soft) 100%); border: 1px solid var(--wz-border); border-radius: 10px; padding: 14px 16px; }
.wz-kpi .v { font-size: 24px; font-weight: 600; color: var(--wz-accent-strong); letter-spacing: -0.01em; font-variant-numeric: tabular-nums; }
.wz-kpi .l { font-size: 11.5px; color: var(--wz-muted); margin-top: 2px; text-transform: uppercase; letter-spacing: 0.03em; font-weight: 500; }
.wz-kpi.ok .v { color: var(--wz-ok); }
.wz-kpi.warn .v { color: var(--wz-warn); }
.wz-kpi.err .v { color: var(--wz-err); }

.wz-tbl-wrap { margin-top: 14px; border: 1px solid var(--wz-border); border-radius: 10px; overflow: hidden; background: #FFFFFF; }
.wz-tbl-wrap header { background: var(--wz-surface-soft); padding: 10px 14px; font-size: 13px; font-weight: 600; color: var(--wz-ink); display: flex; justify-content: space-between; align-items: center; }
.wz-tbl-wrap .wz-count { color: var(--wz-muted); font-size: 12px; font-weight: 400; }
.wz-tbl-wrap .wz-scroll { max-height: 360px; overflow: auto; }
.wz-tbl-wrap table { width: 100%; border-collapse: collapse; font-size: 12.5px; }
.wz-tbl-wrap th { background: #EEF5FC; color: var(--wz-ink); font-weight: 500; text-align: left; padding: 8px 10px; border-bottom: 1px solid var(--wz-border); position: sticky; top: 0; }
.wz-tbl-wrap td { padding: 7px 10px; border-bottom: 1px solid #EEF3F8; color: var(--wz-ink); }
.wz-tbl-wrap tr:nth-child(even) td { background: #FAFCFE; }
.wz-tbl-wrap td.num { text-align: right; font-variant-numeric: tabular-nums; font-family: 'IBM Plex Mono', monospace; }
.wz-tbl-wrap td.delta-pos { color: var(--wz-ok); font-weight: 500; }
.wz-tbl-wrap td.delta-neg { color: var(--wz-err); font-weight: 500; }

.wz-banner { padding: 12px 14px; border-radius: 10px; margin: 12px 0; font-size: 13px; display: flex; align-items: flex-start; gap: 10px; }
.wz-banner.info { background: var(--wz-pastel); color: var(--wz-accent-strong); border: 1px solid var(--wz-accent-soft); }
.wz-banner.warn { background: var(--wz-warn-soft); color: var(--wz-warn); border: 1px solid #F5D6A0; }
.wz-banner.err { background: var(--wz-err-soft); color: var(--wz-err); border: 1px solid #F3B8B3; }
.wz-banner.ok { background: var(--wz-ok-soft); color: var(--wz-ok); border: 1px solid #B7DDBE; }

.wz-signoff { background: var(--wz-surface-soft); border: 1px solid var(--wz-border); border-radius: 10px; padding: 14px 16px; margin-top: 16px; font-family: 'IBM Plex Mono', monospace; font-size: 12px; color: var(--wz-ink-soft); line-height: 1.8; }
.wz-signoff b { color: var(--wz-ink); }

.wz-dl-bar { display: flex; flex-wrap: wrap; gap: 8px; margin-top: 14px; padding: 14px; background: linear-gradient(135deg, var(--wz-pastel) 0%, #F1F7FD 100%); border: 1px solid var(--wz-accent-soft); border-radius: 10px; }
.wz-dl-bar .wz-lbl { font-weight: 600; color: var(--wz-accent-strong); margin-right: 6px; display: flex; align-items: center; }
.wz-toast { position: fixed; right: 20px; bottom: 20px; padding: 10px 16px; background: var(--wz-ink); color: #FFFFFF; border-radius: 8px; font-size: 13px; box-shadow: 0 6px 20px rgba(15,42,68,0.20); opacity: 0; transform: translateY(8px); transition: opacity 180ms, transform 180ms; pointer-events: none; z-index: 1000; max-width: 360px; }
.wz-toast.show { opacity: 1; transform: translateY(0); }
.wz-toast.err { background: var(--wz-err); }
.wz-toast.ok { background: var(--wz-ok); }
.wz-toast.warn { background: var(--wz-warn); color: #FFFFFF; }

.wz-tagrow { display: flex; flex-wrap: wrap; gap: 6px; margin-top: 6px; }
.wz-tag { display: inline-flex; align-items: center; gap: 4px; padding: 3px 10px; background: var(--wz-pastel); color: var(--wz-accent-strong); border-radius: 999px; font-size: 11.5px; border: 1px solid var(--wz-accent-soft); font-weight: 500; }
.wz-tag.removable { padding-right: 4px; }
.wz-tag .wz-x { background: transparent; border: none; cursor: pointer; color: var(--wz-accent-strong); font-size: 14px; line-height: 1; padding: 0 4px; }
.wz-tag .wz-x:hover { color: var(--wz-err); }

.wz-switch { display: inline-flex; align-items: center; gap: 8px; cursor: pointer; user-select: none; }
.wz-switch input { display: none; }
.wz-switch .wz-track { width: 36px; height: 20px; background: var(--wz-pastel-2); border-radius: 999px; position: relative; transition: background 140ms; }
.wz-switch .wz-track::after { content: ''; position: absolute; top: 2px; left: 2px; width: 16px; height: 16px; background: #FFFFFF; border-radius: 50%; transition: left 140ms; box-shadow: 0 1px 2px rgba(0,0,0,0.2); }
.wz-switch input:checked + .wz-track { background: var(--wz-accent); }
.wz-switch input:checked + .wz-track::after { left: 18px; }
.wz-switch .wz-switch-lbl { font-size: 13px; color: var(--wz-ink-soft); }

@media (max-width: 720px) {
  .wz-root { padding: 12px 10px 40px; }
  .wz-panel { padding: 16px 14px; }
  .wz-step { min-width: 140px; }
  .wz-step .wz-tag { font-size: 10px; }
  .wz-step .wz-lbl { font-size: 12.5px; }
}
"""

# ---------------------------------------------------------------------------
# CDN libraries for client-side download (primary path)
# ---------------------------------------------------------------------------

CDN_SCRIPTS = """
<script defer src="https://cdn.jsdelivr.net/npm/xlsx@0.18.5/dist/xlsx.full.min.js" referrerpolicy="no-referrer"></script>
<script defer src="https://cdn.jsdelivr.net/npm/jszip@3.10.1/dist/jszip.min.js" referrerpolicy="no-referrer"></script>
<script>
  window.__WZ_CDN__ = {
    mammoth: 'https://cdn.jsdelivr.net/npm/mammoth@1.6.0/mammoth.browser.min.js',
    pdfjs:   'https://cdn.jsdelivr.net/npm/pdfjs-dist@3.11.174/build/pdf.min.js',
    pdfjsWorker: 'https://cdn.jsdelivr.net/npm/pdfjs-dist@3.11.174/build/pdf.worker.min.js',
    jspdf:   'https://cdn.jsdelivr.net/npm/jspdf@2.5.1/dist/jspdf.umd.min.js',
    autotable: 'https://cdn.jsdelivr.net/npm/jspdf-autotable@3.8.2/dist/jspdf.plugin.autotable.min.js',
    docx:    'https://cdn.jsdelivr.net/npm/docx@8.5.0/build/index.umd.min.js',
    htmlDocx: 'https://cdn.jsdelivr.net/npm/html-docx-js@0.3.1/dist/html-docx.js'
  };
</script>
"""

# ---------------------------------------------------------------------------
# Streaming observer: watches parent chat for @@@RECON-START/END and renders
# ---------------------------------------------------------------------------

OBSERVER_SCRIPT = r"""
<script>
(function(){
  const START = '@@@RECON-START';
  const END = '@@@RECON-END';
  const root = document.getElementById('ibm-render');
  const loader = document.getElementById('ibm-loader');
  const status = document.getElementById('ibm-status');
  const NAVY = '__IBM_NAVY__', BLUE = '__IBM_BLUE__';
  const SERVER_FALLBACK = window.__IBM_SERVER_FALLBACK__ || {};

  function toast(msg, isErr){
    const t = document.getElementById('ibm-toast');
    if(!t) return;
    t.textContent = msg;
    t.classList.toggle('err', !!isErr);
    t.classList.add('show');
    setTimeout(()=>t.classList.remove('show'), 2600);
  }

  function esc(s){ return String(s==null?'':s).replace(/[&<>"']/g, c=>({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}[c])); }
  function fmt(n){
    if(n==null||n==='') return '';
    if(typeof n==='number') return n.toLocaleString(undefined,{maximumFractionDigits:2,minimumFractionDigits:(Math.abs(n)%1>0?2:0)});
    const f=parseFloat(String(n).replace(/,/g,''));
    return isFinite(f)?f.toLocaleString(undefined,{maximumFractionDigits:2,minimumFractionDigits:(Math.abs(f)%1>0?2:0)}):String(n);
  }

  function renderHeader(meta){
    const chips = [];
    if(meta.sector) chips.push(`<span class="ibm-chip primary">${esc(meta.sector)}</span>`);
    if(meta.region) chips.push(`<span class="ibm-chip">${esc(meta.region)}</span>`);
    (meta.regulations||[]).forEach(r=>chips.push(`<span class="ibm-regtag">${esc(r)}</span>`));
    if(meta.asOf) chips.push(`<span class="ibm-chip">as of ${esc(meta.asOf)}</span>`);
    if(meta.redacted) chips.push(`<span class="ibm-chip warn">PII redacted</span>`);
    return `<div class="ibm-subbar">${chips.join('')}</div>`;
  }

  function renderKpis(s){
    const items = [
      {v:s.total_left||0, l:'Records (Left)'},
      {v:s.total_right||0, l:'Records (Right)'},
      {v:s.matched||0, l:'Matched', cls:'ok'},
      {v:s.variance||0, l:'Variance', cls:'warn'},
      {v:s.unmatched_left||0, l:'Missing in Right', cls:(s.unmatched_left>0?'danger':'')},
      {v:s.unmatched_right||0, l:'Missing in Left', cls:(s.unmatched_right>0?'danger':'')},
    ];
    if(s.match_rate!=null) items.push({v:s.match_rate+'%', l:'Match rate', cls:(s.match_rate>=95?'ok':s.match_rate>=80?'warn':'danger')});
    return `<div class="ibm-card"><header>Reconciliation KPIs</header><div class="body"><div class="ibm-kpis">${
      items.map(i=>`<div class="ibm-kpi ${i.cls||''}"><div class="v">${esc(fmt(i.v))}</div><div class="l">${esc(i.l)}</div></div>`).join('')
    }</div></div></div>`;
  }

  function renderTable(title, cols, rows, extraClass){
    if(!rows||!rows.length) return `<div class="ibm-card"><header>${esc(title)} <span class="ibm-chip success">0</span></header><div class="body"><em>None.</em></div></div>`;
    const th = cols.map(c=>`<th class="${c.num?'num':''}">${esc(c.label||c.key)}</th>`).join('');
    const tr = rows.map(r=>{
      const tds = cols.map(c=>{
        let v = r[c.key];
        let cls = c.num?'num':'';
        if(c.key==='delta'||c.key==='variance'){
          const nv = parseFloat(v);
          if(isFinite(nv)) cls += (nv<0?' delta-neg':(nv>0?' delta-pos':''));
        }
        if(c.keyField) cls += ' key';
        return `<td class="${cls}">${esc(c.num?fmt(v):v)}</td>`;
      }).join('');
      return `<tr>${tds}</tr>`;
    }).join('');
    return `<div class="ibm-card ${extraClass||''}"><header>${esc(title)} <span class="ibm-chip">${rows.length} row${rows.length===1?'':'s'}</span></header><div class="body"><div class="ibm-tab-scroll"><table class="ibm-table"><thead><tr>${th}</tr></thead><tbody>${tr}</tbody></table></div></div></div>`;
  }

  function renderNarrative(items){
    if(!items||!items.length) return '';
    const body = items.map(n=>{
      const tag = n.regulation?`<span class="ibm-regtag">${esc(n.regulation)}</span>`:'';
      const sev = n.severity==='high'?'danger':(n.severity==='medium'?'warn':'');
      const chip = sev?`<span class="ibm-chip ${sev}">${esc(n.severity)}</span> `:'';
      return `<p>${chip}${tag}${esc(n.text||'')}</p>`;
    }).join('');
    return `<div class="ibm-section-title"><span class="bar"></span>Variance commentary</div><div class="ibm-narrative">${body}</div>`;
  }

  function renderFooter(meta){
    return `<div class="ibm-footer"><span>IBM Consulting Advantage — AI Reconciliation</span><span>Build __BUILD__ · Run ${esc(meta.runId||'')}</span></div>`;
  }

  // --- Payload state (populated when RECON block arrives) ---
  let PAYLOAD = null;

  function renderAll(p){
    PAYLOAD = p;
    const meta = p.meta||{};
    const stats = p.stats||{};
    const parts = [];
    parts.push(renderHeader(meta));
    parts.push(renderKpis(stats));

    const matched = p.matched||[];
    const variance = p.variance||[];
    const ul = p.unmatched_left||[];
    const ur = p.unmatched_right||[];

    if(variance.length){
      const cols = p.variance_columns || inferCols(variance[0]);
      parts.push(`<div class="ibm-section-title"><span class="bar"></span>Variances — amount / field mismatches</div>`);
      parts.push(renderTable('Variance detail', cols, variance));
    }
    if(ul.length){
      const cols = p.unmatched_columns || inferCols(ul[0]);
      parts.push(`<div class="ibm-section-title"><span class="bar"></span>Missing in Right side</div>`);
      parts.push(renderTable('Missing in Right', cols, ul));
    }
    if(ur.length){
      const cols = p.unmatched_columns || inferCols(ur[0]);
      parts.push(`<div class="ibm-section-title"><span class="bar"></span>Missing in Left side</div>`);
      parts.push(renderTable('Missing in Left', cols, ur));
    }
    if(matched.length){
      const cols = p.matched_columns || inferCols(matched[0]);
      parts.push(`<div class="ibm-section-title"><span class="bar"></span>Matched records (preview)</div>`);
      parts.push(renderTable('Matched', cols, matched.slice(0,25)));
    }
    parts.push(renderNarrative(p.narrative));
    parts.push(renderFooter(meta));

    root.innerHTML = parts.join('\n');
    if(loader) loader.remove();
    status.textContent = 'Reconciliation complete';
    updateDownloadButtons(true);
  }

  function inferCols(r){
    return Object.keys(r).map(k=>({key:k, label:k, num:/amount|qty|quantity|price|value|delta|variance|paid|allowed|charge|premium|reserve/i.test(k)}));
  }

  // --- Download helpers ---
  function updateDownloadButtons(enabled){
    ['ibm-dl-xlsx','ibm-dl-docx','ibm-dl-pptx'].forEach(id=>{
      const b = document.getElementById(id);
      if(b) b.disabled = !enabled;
    });
  }

  function triggerDataUri(name, dataUri){
    const a = document.createElement('a');
    a.href = dataUri;
    a.download = name;
    document.body.appendChild(a);
    a.click();
    setTimeout(()=>a.remove(), 200);
  }

  function fallbackDownload(fmtName){
    const payload = SERVER_FALLBACK[fmtName];
    if(!payload){ toast('Server fallback not available', true); return; }
    triggerDataUri(payload.filename, payload.dataUri);
    toast('Downloaded via ooXML fallback');
  }

  window.ibmDownload = function(fmtName){
    if(!PAYLOAD){ toast('Reconciliation payload not ready', true); return; }
    const meta = PAYLOAD.meta||{};
    const baseName = (meta.outputName || 'reconciliation') + '_' + (meta.asOf||'').replace(/[^0-9A-Za-z]/g,'');
    try{
      if(fmtName==='xlsx'){
        if(typeof XLSX==='undefined') throw new Error('XLSX CDN unavailable');
        const wb = XLSX.utils.book_new();
        const summary = [['IBM Consulting Advantage — Reconciliation'],[],['Sector', meta.sector||''],['Region', meta.region||''],['As of', meta.asOf||''],['Run ID', meta.runId||''],[],['KPI','Value']];
        Object.entries(PAYLOAD.stats||{}).forEach(([k,v])=>summary.push([k, v]));
        XLSX.utils.book_append_sheet(wb, XLSX.utils.aoa_to_sheet(summary), 'Summary');
        const sections = [['Variance', PAYLOAD.variance||[]],['Missing_in_Right', PAYLOAD.unmatched_left||[]],['Missing_in_Left', PAYLOAD.unmatched_right||[]],['Matched', PAYLOAD.matched||[]]];
        sections.forEach(([name, rows])=>{
          if(!rows.length) return;
          XLSX.utils.book_append_sheet(wb, XLSX.utils.json_to_sheet(rows), name);
        });
        XLSX.writeFile(wb, baseName+'.xlsx');
        toast('XLSX downloaded via CDN');
      } else if(fmtName==='docx'){
        if(typeof window.htmlDocx==='undefined') throw new Error('html-docx CDN unavailable');
        const shell = buildDocxHtml(PAYLOAD);
        const blob = window.htmlDocx.asBlob(shell);
        const url = URL.createObjectURL(blob);
        const a = document.createElement('a'); a.href=url; a.download=baseName+'.docx'; document.body.appendChild(a); a.click();
        setTimeout(()=>{URL.revokeObjectURL(url); a.remove();},400);
        toast('DOCX downloaded via CDN');
      } else if(fmtName==='pptx'){
        if(typeof PptxGenJS==='undefined') throw new Error('PptxGenJS CDN unavailable');
        buildPptx(PAYLOAD).then(()=>toast('PPTX downloaded via CDN'));
      }
    } catch(e){
      console.warn('CDN path failed, using ooXML fallback:', e.message);
      fallbackDownload(fmtName);
    }
  };

  function buildDocxHtml(p){
    const meta=p.meta||{}, s=p.stats||{};
    const tbl = (rows, caption)=>{
      if(!rows.length) return '';
      const keys = Object.keys(rows[0]);
      const th = keys.map(k=>`<th style="background:${NAVY};color:#fff;padding:6px;border:1px solid #888;">${k}</th>`).join('');
      const tr = rows.map(r=>`<tr>${keys.map(k=>`<td style="padding:6px;border:1px solid #ccc;">${(r[k]==null?'':String(r[k]).replace(/[<>&]/g,''))}</td>`).join('')}</tr>`).join('');
      return `<h3 style="color:${NAVY};">${caption} (${rows.length})</h3><table style="border-collapse:collapse;width:100%;font-size:11px;"><thead><tr>${th}</tr></thead><tbody>${tr}</tbody></table>`;
    };
    const narr = (p.narrative||[]).map(n=>`<p style="border-left:3px solid ${BLUE};padding-left:10px;margin:6px 0;"><strong>${n.regulation||''}</strong> ${n.text||''}</p>`).join('');
    const kpis = Object.entries(s).map(([k,v])=>`<li><strong>${k}:</strong> ${v}</li>`).join('');
    return `<!DOCTYPE html><html><head><meta charset="utf-8"><title>Reconciliation</title></head><body style="font-family:'Segoe UI',Arial,sans-serif;color:#161616;">
    <div style="background:${NAVY};color:#fff;padding:16px;"><h1 style="margin:0;">IBM Consulting Advantage — AI Reconciliation</h1>
    <div style="opacity:.85;font-size:12px;">${meta.sector||''} · ${meta.region||''} · ${meta.asOf||''}</div></div>
    <h2 style="color:${NAVY};">Summary</h2><ul>${kpis}</ul>
    ${narr?`<h2 style="color:${NAVY};">Commentary</h2>${narr}`:''}
    ${tbl(p.variance||[], 'Variances')}
    ${tbl(p.unmatched_left||[], 'Missing in Right')}
    ${tbl(p.unmatched_right||[], 'Missing in Left')}
    ${tbl((p.matched||[]).slice(0,50), 'Matched (sample)')}
    </body></html>`;
  }

  async function buildPptx(p){
    const pres = new PptxGenJS();
    pres.layout = 'LAYOUT_WIDE';
    const meta=p.meta||{}, s=p.stats||{};
    const navy = NAVY.replace('#',''), blue=BLUE.replace('#','');
    // Title slide
    let sl = pres.addSlide(); sl.background={color:navy};
    sl.addText('IBM Consulting Advantage', {x:0.5,y:0.4,w:12,h:0.5,fontSize:18,color:'FFFFFF',bold:true});
    sl.addText('AI Reconciliation Report', {x:0.5,y:1.2,w:12,h:1.0,fontSize:40,color:'FFFFFF',bold:true});
    sl.addText(`${meta.sector||''} · ${meta.region||''} · ${meta.asOf||''}`, {x:0.5,y:2.4,w:12,h:0.5,fontSize:18,color:'D0E2FF'});
    // KPI slide
    sl = pres.addSlide(); sl.background={color:'FFFFFF'};
    sl.addShape(pres.ShapeType.rect,{x:0,y:0,w:13.33,h:0.6,fill:{color:navy}});
    sl.addText('Reconciliation KPIs', {x:0.3,y:0.05,w:12,h:0.5,fontSize:20,color:'FFFFFF',bold:true});
    const kpiRows = Object.entries(s);
    const perRow = 4, cw = 2.9, ch = 1.1;
    kpiRows.forEach(([k,v],i)=>{
      const col = i%perRow, row = Math.floor(i/perRow);
      const x = 0.5 + col*(cw+0.1), y = 1.0 + row*(ch+0.2);
      sl.addShape(pres.ShapeType.rect,{x,y,w:cw,h:ch,fill:{color:'EDF5FF'},line:{color:'D0E2FF',width:1}});
      sl.addText(String(v), {x:x,y:y,w:cw,h:ch*0.55,fontSize:22,color:navy,bold:true,align:'center',valign:'middle'});
      sl.addText(k, {x:x,y:y+ch*0.55,w:cw,h:ch*0.45,fontSize:10,color:'525252',align:'center'});
    });
    // Table slides for each section
    const sections = [['Variances', p.variance||[]],['Missing in Right', p.unmatched_left||[]],['Missing in Left', p.unmatched_right||[]]];
    sections.forEach(([title, rows])=>{
      if(!rows.length) return;
      sl = pres.addSlide(); sl.background={color:'FFFFFF'};
      sl.addShape(pres.ShapeType.rect,{x:0,y:0,w:13.33,h:0.6,fill:{color:navy}});
      sl.addText(title+` (${rows.length})`, {x:0.3,y:0.05,w:12,h:0.5,fontSize:20,color:'FFFFFF',bold:true});
      const keys = Object.keys(rows[0]).slice(0,6);
      const header = keys.map(k=>({text:k,options:{bold:true,color:'FFFFFF',fill:{color:navy}}}));
      const data = [header].concat(rows.slice(0,15).map(r=>keys.map(k=>({text:String(r[k]==null?'':r[k])}))));
      sl.addTable(data, {x:0.3,y:0.8,w:12.7,fontSize:9,border:{type:'solid',pt:0.5,color:'CCCCCC'}});
    });
    // Narrative slide
    if((p.narrative||[]).length){
      sl = pres.addSlide(); sl.background={color:'FFFFFF'};
      sl.addShape(pres.ShapeType.rect,{x:0,y:0,w:13.33,h:0.6,fill:{color:navy}});
      sl.addText('Variance Commentary', {x:0.3,y:0.05,w:12,h:0.5,fontSize:20,color:'FFFFFF',bold:true});
      const bullets = (p.narrative||[]).slice(0,8).map(n=>({text:(n.regulation?`[${n.regulation}] `:'')+(n.text||''),options:{bullet:true,fontSize:14,color:'161616'}}));
      sl.addText(bullets,{x:0.5,y:0.9,w:12.3,h:6.2});
    }
    const name = (meta.outputName || 'reconciliation') + '_' + (meta.asOf||'').replace(/[^0-9A-Za-z]/g,'') + '.pptx';
    await pres.writeFile({fileName:name});
  }

  // --- Observer: tail parent chat for @@@RECON block ---
  function readParentText(){
    try{
      const doc = window.parent && window.parent.document;
      if(!doc) return null;
      // Find the LAST message bubble that contains our markers.
      // Fall back to concatenated textContent of all message-content nodes.
      const nodes = doc.querySelectorAll('.chat-assistant .prose, .message-content, [data-message-role="assistant"]');
      if(nodes && nodes.length){
        for(let i=nodes.length-1;i>=0;i--){
          const t = nodes[i].innerText || nodes[i].textContent || '';
          if(t.indexOf(START)>=0) return t;
        }
        return nodes[nodes.length-1].innerText || '';
      }
      return doc.body ? doc.body.innerText : '';
    }catch(e){ return null; }
  }

  function extractBlock(text){
    if(!text) return null;
    const i = text.indexOf(START); if(i<0) return null;
    const j = text.indexOf(END, i); if(j<0) return null;
    return text.slice(i+START.length, j).trim();
  }

  function tryParse(raw){
    if(!raw) return null;
    // Strip accidental code fences
    raw = raw.replace(/^```(?:json)?/i,'').replace(/```$/,'').trim();
    try { return JSON.parse(raw); } catch(e){}
    // Extract largest {...} blob
    const m = raw.match(/\{[\s\S]*\}/);
    if(m){ try{ return JSON.parse(m[0]); }catch(e){} }
    return null;
  }

  let done = false;
  function tick(){
    if(done) return;
    const txt = readParentText();
    const block = extractBlock(txt||'');
    if(block){
      const obj = tryParse(block);
      if(obj && obj.stats){
        done = true;
        renderAll(obj);
        try{ hideMarkersInParent(); }catch(e){}
      }
    }
  }

  function hideMarkersInParent(){
    try{
      const doc = window.parent.document;
      const style = doc.createElement('style');
      style.textContent = `.ibm-recon-hidden{display:none!important;}`;
      doc.head.appendChild(style);
      // Not a perfect hider, but most OWUI themes collapse fenced text — we leave rendered.
    }catch(e){}
  }

  // Also accept payload posted directly from Python (inline server-rendered payload)
  if(window.__IBM_RECON_PAYLOAD__){
    try { renderAll(window.__IBM_RECON_PAYLOAD__); } catch(e){ console.error(e); }
  }

  const iv = setInterval(()=>{ if(done){ clearInterval(iv); return; } tick(); }, 600);
  setTimeout(()=>{ if(!done){ if(status) status.textContent = 'Awaiting reconciliation stream…'; } }, 3000);

  // Wire up buttons
  document.addEventListener('click', (e)=>{
    const t = e.target.closest('[data-ibm-dl]');
    if(!t) return;
    window.ibmDownload(t.getAttribute('data-ibm-dl'));
  });

  // Height auto-size
  function autoHeight(){
    try{
      const h = Math.max(document.documentElement.scrollHeight, document.body.scrollHeight);
      window.parent.postMessage({type:'iframe:resize', height:h}, '*');
    }catch(e){}
  }
  setInterval(autoHeight, 700);
  window.addEventListener('load', autoHeight);
})();
</script>
"""

# ---------------------------------------------------------------------------
# Wizard script: upload, parse, reconcile, and download — all client-side
# ---------------------------------------------------------------------------

WIZARD_SCRIPT = r"""
<script>
(function(){
  'use strict';
  const CFG = window.__WZ_CFG__ || {};
  const SECTORS = CFG.sectors || {};
  const REGIONS = CFG.regions || ['USA', 'European Union', 'United Kingdom', 'Global'];
  const BUILD = CFG.build || '';

  const $ = (sel, el) => (el || document).querySelector(sel);
  const $$ = (sel, el) => Array.from((el || document).querySelectorAll(sel));
  const el = (tag, attrs, children) => {
    const n = document.createElement(tag);
    if (attrs) for (const k in attrs) {
      if (k === 'class') n.className = attrs[k];
      else if (k === 'style') n.style.cssText = attrs[k];
      else if (k.startsWith('on') && typeof attrs[k] === 'function') n.addEventListener(k.slice(2), attrs[k]);
      else if (k === 'html') n.innerHTML = attrs[k];
      else if (attrs[k] != null) n.setAttribute(k, attrs[k]);
    }
    if (children) (Array.isArray(children) ? children : [children]).forEach(c => { if (c == null) return; n.appendChild(typeof c === 'string' ? document.createTextNode(c) : c); });
    return n;
  };

  const STEPS = [
    { id: 1, tag: 'Step 1', label: 'Scope & Jurisdiction' },
    { id: 2, tag: 'Step 2', label: 'Upload Datasets' },
    { id: 3, tag: 'Step 3', label: 'Matching Rules' },
    { id: 4, tag: 'Step 4', label: 'Run Reconciliation' },
    { id: 5, tag: 'Final',  label: 'Review & Download' },
  ];

  const state = {
    step: 1,
    scope: {
      region: 'USA',
      sector: 'Banking',
      asOf: new Date().toISOString().slice(0, 10),
      runId: genRunId(),
      operator: '',
      outputName: 'reconciliation',
    },
    left:  { files: [], tables: [], selectedTable: 0, records: [], sha256: null, status: 'empty' },
    right: { files: [], tables: [], selectedTable: 0, records: [], sha256: null, status: 'empty' },
    rules: {
      keyFields: [],
      amountFields: [],
      tolerance: 0.01,
      tolerancePct: 0,
      dateTolerance: 1,
      currencyMode: 'strict',
      piiRedact: false,
    },
    result: null,
    cdnStatus: {
      XLSX: typeof XLSX !== 'undefined',
      JSZip: typeof JSZip !== 'undefined',
      mammoth: typeof mammoth !== 'undefined',
      pdfjsLib: typeof pdfjsLib !== 'undefined',
      jspdf: typeof (window.jspdf || {}).jsPDF !== 'undefined',
      docx: typeof (window.docx || {}).Document !== 'undefined',
      htmlDocx: typeof window.htmlDocx !== 'undefined',
    },
  };

  // --- Lazy script loader ----------------------------------------------------
  // Only SheetJS + JSZip are loaded eagerly (deferred) for fast first paint.
  // Heavier libs load on demand when the user actually needs them.
  const _loaded = {};
  function loadScript(url) {
    if (_loaded[url]) return _loaded[url];
    _loaded[url] = new Promise((resolve, reject) => {
      const s = document.createElement('script');
      s.src = url; s.async = true; s.referrerPolicy = 'no-referrer';
      s.onload = () => resolve(true);
      s.onerror = () => { delete _loaded[url]; reject(new Error('Failed to load ' + url)); };
      document.head.appendChild(s);
    });
    return _loaded[url];
  }
  async function ensureSheetJS() {
    if (typeof XLSX !== 'undefined') { state.cdnStatus.XLSX = true; return true; }
    // SheetJS is deferred-eager; if still missing after short wait, give up
    for (let i = 0; i < 20; i++) { if (typeof XLSX !== 'undefined') { state.cdnStatus.XLSX = true; return true; } await sleep(120); }
    return false;
  }
  async function ensureJSZip() {
    if (typeof JSZip !== 'undefined') { state.cdnStatus.JSZip = true; return true; }
    for (let i = 0; i < 20; i++) { if (typeof JSZip !== 'undefined') { state.cdnStatus.JSZip = true; return true; } await sleep(120); }
    return false;
  }
  async function ensureMammoth() {
    if (typeof mammoth !== 'undefined') { state.cdnStatus.mammoth = true; return true; }
    try { await loadScript(window.__WZ_CDN__.mammoth); state.cdnStatus.mammoth = typeof mammoth !== 'undefined'; if (!state.cdnStatus.mammoth) noteCdnBlocked('DOCX (mammoth)'); return state.cdnStatus.mammoth; } catch (e) { noteCdnBlocked('DOCX (mammoth)'); return false; }
  }
  async function ensurePdfJs() {
    if (typeof pdfjsLib !== 'undefined') { state.cdnStatus.pdfjsLib = true; return true; }
    try {
      await loadScript(window.__WZ_CDN__.pdfjs);
      if (typeof pdfjsLib !== 'undefined') {
        try { pdfjsLib.GlobalWorkerOptions.workerSrc = window.__WZ_CDN__.pdfjsWorker; } catch (e) {}
      }
      state.cdnStatus.pdfjsLib = typeof pdfjsLib !== 'undefined';
      if (!state.cdnStatus.pdfjsLib) noteCdnBlocked('PDF (pdf.js)');
      return state.cdnStatus.pdfjsLib;
    } catch (e) { noteCdnBlocked('PDF (pdf.js)'); return false; }
  }
  async function ensureJsPDF() {
    if (((window.jspdf || {}).jsPDF)) { state.cdnStatus.jspdf = true; return true; }
    try { await loadScript(window.__WZ_CDN__.jspdf); await loadScript(window.__WZ_CDN__.autotable); state.cdnStatus.jspdf = !!((window.jspdf || {}).jsPDF); if (!state.cdnStatus.jspdf) noteCdnBlocked('PDF output (jsPDF)'); return state.cdnStatus.jspdf; } catch (e) { noteCdnBlocked('PDF output (jsPDF)'); return false; }
  }
  async function ensureDocxJs() {
    if (((window.docx || {}).Document)) { state.cdnStatus.docx = true; return true; }
    try { await loadScript(window.__WZ_CDN__.docx); state.cdnStatus.docx = !!((window.docx || {}).Document); return state.cdnStatus.docx; } catch (e) { return false; }
  }
  async function ensureHtmlDocx() {
    if (window.htmlDocx) { state.cdnStatus.htmlDocx = true; return true; }
    try { await loadScript(window.__WZ_CDN__.htmlDocx); state.cdnStatus.htmlDocx = !!window.htmlDocx; if (!state.cdnStatus.htmlDocx && !state.cdnStatus.docx) noteCdnBlocked('DOCX output'); return state.cdnStatus.htmlDocx; } catch (e) { if (!state.cdnStatus.docx) noteCdnBlocked('DOCX output'); return false; }
  }

  function genRunId() {
    const r = crypto && crypto.getRandomValues ? crypto.getRandomValues(new Uint8Array(4)) : [Math.random()*255, Math.random()*255, Math.random()*255, Math.random()*255];
    return Array.from(r).map(b => (b & 255).toString(16).padStart(2, '0')).join('').slice(0, 8);
  }

  async function sha256Hex(buf) {
    try {
      const hash = await crypto.subtle.digest('SHA-256', buf);
      return Array.from(new Uint8Array(hash)).map(b => b.toString(16).padStart(2, '0')).join('');
    } catch (e) { return null; }
  }

  function toast(msg, kind) {
    const t = $('#wz-toast'); if (!t) return;
    t.textContent = msg;
    t.className = 'wz-toast show' + (kind ? ' ' + kind : '');
    clearTimeout(toast._t);
    toast._t = setTimeout(() => t.classList.remove('show'), 3200);
  }

  function esc(s) { return String(s == null ? '' : s).replace(/[&<>"']/g, c => ({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}[c])); }
  function fmtNum(v) {
    if (v == null || v === '') return '';
    if (typeof v === 'number') return v.toLocaleString(undefined, { maximumFractionDigits: 2 });
    const n = parseFloat(String(v).replace(/,/g, ''));
    return isFinite(n) ? n.toLocaleString(undefined, { maximumFractionDigits: 2 }) : String(v);
  }
  function isAmountKey(k) { return /amount|qty|quantity|price|value|delta|variance|paid|allowed|charge|premium|reserve|total|net|gross|tax|balance|usage|kwh|billed|obligated|disbursed/i.test(k); }
  function isIdKey(k) { return /id$|^id|ref$|number$|#$|code$|isin|lei|claim|trade|policy|account|subject|voucher|invoice|po|order|shipment|meter|subscriber/i.test(k); }

  function renderStepper() {
    const container = $('#wz-stepper'); if (!container) return;
    container.innerHTML = '';
    STEPS.forEach(s => {
      const done = s.id < state.step;
      const active = s.id === state.step;
      const node = el('div', { class: 'wz-step' + (active ? ' active' : '') + (done ? ' done' : ''), onclick: () => { if (s.id < state.step) goto(s.id); } }, [
        el('div', { class: 'wz-bullet' }, done ? '' : String(s.id)),
        el('div', { class: 'wz-txt' }, [
          el('span', { class: 'wz-tag' }, s.tag),
          el('span', { class: 'wz-lbl' }, s.label),
        ]),
      ]);
      container.appendChild(node);
    });
  }

  function goto(step) {
    state.step = Math.max(1, Math.min(5, step));
    renderStepper();
    renderBody();
  }

  function renderBody() {
    const body = $('#wz-body'); if (!body) return;
    body.innerHTML = '';
    if (state.step === 1) body.appendChild(renderStep1());
    if (state.step === 2) body.appendChild(renderStep2());
    if (state.step === 3) body.appendChild(renderStep3());
    if (state.step === 4) body.appendChild(renderStep4());
    if (state.step === 5) body.appendChild(renderStep5());
    autoHeight();
  }

  function renderCdnBanner() {
    // Heavy libs load on demand — only show a banner if a lazy load has actually failed.
    if (!state.cdnBlocked || !state.cdnBlocked.length) return null;
    return el('div', { class: 'wz-banner warn', html:
      '<div><b>CDN blocked for:</b> <i>' + esc(state.cdnBlocked.join(', ')) + '</i>. ' +
      'CSV, TSV, and JSON still work natively; you can also download CSV results. For XLSX / DOCX / PPTX / PDF inputs or outputs, attach the files in chat and ask the assistant to reconcile — the server produces IBM-branded ooXML as a fallback.</div>'
    });
  }
  function noteCdnBlocked(name) {
    state.cdnBlocked = state.cdnBlocked || [];
    if (state.cdnBlocked.indexOf(name) < 0) state.cdnBlocked.push(name);
  }

  // --- Step 1: Scope ---------------------------------------------------------
  function renderStep1() {
    const panel = el('div', { class: 'wz-panel' });
    panel.appendChild(el('h2', {}, STEPS[0].tag + ' — ' + STEPS[0].label));
    panel.appendChild(el('p', { class: 'wz-sub' }, 'Choose the regulatory region and industry. Defaults drive key/amount fields and regulation tags for the audit trail.'));
    const banner = renderCdnBanner(); if (banner) panel.appendChild(banner);

    const grid = el('div', { class: 'wz-grid' });
    grid.appendChild(field('Region / Jurisdiction',
      selectInput(REGIONS, state.scope.region, v => { state.scope.region = v; syncDefaults(); }),
      'Regulation tags shown in the audit memo follow this jurisdiction.'
    ));
    grid.appendChild(field('Industry / Sector',
      selectInput(Object.keys(SECTORS), state.scope.sector, v => { state.scope.sector = v; syncDefaults(); }),
      'Select "Other" for sectors not listed. Defaults can be overridden in Step 3.'
    ));
    grid.appendChild(field('Reporting "as of" date',
      inputNode('date', state.scope.asOf, v => state.scope.asOf = v),
      'Shown in the header of every download.'
    ));
    grid.appendChild(field('Run ID (auto)',
      inputNode('text', state.scope.runId, v => state.scope.runId = v, true),
      'Immutable identifier stamped on every output for audit traceability.'
    ));
    grid.appendChild(field('Operator / Analyst',
      inputNode('text', state.scope.operator, v => state.scope.operator = v, false, 'Name shown on sign-off block'),
      'Appears in the sign-off footer of all downloads.'
    ));
    grid.appendChild(field('Output filename base',
      inputNode('text', state.scope.outputName, v => state.scope.outputName = v),
      'File downloads will be named "<base>_<asOf>.<ext>".'
    ));
    panel.appendChild(grid);

    // Regulation tag preview
    const defaults = SECTORS[state.scope.sector] || SECTORS['Other'];
    const regs = (defaults && defaults.regulations && defaults.regulations[state.scope.region]) || [];
    if (regs.length) {
      const tagRow = el('div', { style: 'margin-top:18px;' });
      tagRow.appendChild(el('label', { style: 'font-size:12.5px;font-weight:500;color:var(--wz-ink-soft);' }, 'Regulation tags (preview)'));
      const row = el('div', { class: 'wz-tagrow' });
      regs.forEach(r => row.appendChild(el('span', { class: 'wz-tag' }, r)));
      tagRow.appendChild(row);
      panel.appendChild(tagRow);
    }

    panel.appendChild(renderActions(false, true, 'Next — upload datasets'));
    return panel;
  }

  function syncDefaults() {
    const defaults = SECTORS[state.scope.sector] || SECTORS['Other'] || { key_fields: [], amount_fields: [] };
    if (!state.rules.keyFields.length) state.rules.keyFields = (defaults.key_fields || []).slice();
    if (!state.rules.amountFields.length) state.rules.amountFields = (defaults.amount_fields || []).slice();
    state.rules.piiRedact = /healthcare|pharma/i.test(state.scope.sector);
    renderBody();
  }

  function field(label, inputEl, hint) {
    const f = el('div', { class: 'wz-field' });
    f.appendChild(el('label', {}, label));
    f.appendChild(inputEl);
    if (hint) f.appendChild(el('div', { class: 'wz-hint' }, hint));
    return f;
  }
  function inputNode(type, value, onchange, readonly, placeholder) {
    const i = el('input', { class: 'wz-input', type: type, value: value || '', placeholder: placeholder || '' });
    if (readonly) i.readOnly = true;
    i.addEventListener('input', () => onchange(i.value));
    return i;
  }
  function selectInput(options, value, onchange) {
    const s = el('select', { class: 'wz-select' });
    options.forEach(o => {
      const opt = el('option', { value: o }, o);
      if (o === value) opt.selected = true;
      s.appendChild(opt);
    });
    s.addEventListener('change', () => onchange(s.value));
    return s;
  }

  // --- Step 2: Upload --------------------------------------------------------
  function renderStep2() {
    const panel = el('div', { class: 'wz-panel' });
    panel.appendChild(el('h2', {}, STEPS[1].tag + ' — ' + STEPS[1].label));
    panel.appendChild(el('p', { class: 'wz-sub' }, 'Drop the two datasets you want to compare. Supported: CSV, TSV, XLSX, XLS, DOCX, PPTX, PDF, JSON. Each side accepts multiple files — we auto-detect tables.'));
    const banner = renderCdnBanner(); if (banner) panel.appendChild(banner);

    const grid = el('div', { class: 'wz-grid' });
    grid.appendChild(renderDropZone('left', 'Left dataset', 'e.g. GL, blotter, policy ledger, 837 claims, EDC export, PO ledger'));
    grid.appendChild(renderDropZone('right', 'Right dataset', 'e.g. custodian, clearing, claims feed, 835 remittance, CRO export, GRN'));
    panel.appendChild(grid);

    // Preview
    ['left', 'right'].forEach(side => {
      const s = state[side];
      if (s.tables.length) panel.appendChild(renderTablePreview(side));
    });

    const canNext = state.left.records.length > 0 && state.right.records.length > 0;
    panel.appendChild(renderActions(true, canNext, canNext ? 'Next — configure matching rules' : 'Upload both datasets to continue'));
    return panel;
  }

  function renderDropZone(side, title, hint) {
    const wrap = el('div', { class: 'wz-field' });
    wrap.appendChild(el('label', {}, title));
    const drop = el('label', { class: 'wz-drop' + (state[side].files.length ? ' wz-has-files' : '') });
    drop.appendChild(el('div', { class: 'wz-icon' }, '+'));
    drop.appendChild(el('div', { class: 'wz-title' }, state[side].files.length ? (state[side].files.length + ' file(s) loaded — click to add more') : 'Drop files here or click to browse'));
    drop.appendChild(el('div', { class: 'wz-sub' }, hint));
    drop.appendChild(el('span', { class: 'wz-browse' }, 'Browse files'));
    const input = el('input', { type: 'file', multiple: 'multiple', accept: '.csv,.tsv,.xlsx,.xls,.docx,.pptx,.pdf,.json,.txt' });
    input.addEventListener('change', e => handleFiles(side, e.target.files));
    drop.appendChild(input);
    drop.addEventListener('dragover', e => { e.preventDefault(); drop.classList.add('wz-hover'); });
    drop.addEventListener('dragleave', () => drop.classList.remove('wz-hover'));
    drop.addEventListener('drop', e => { e.preventDefault(); drop.classList.remove('wz-hover'); handleFiles(side, e.dataTransfer.files); });
    wrap.appendChild(drop);

    if (state[side].files.length) {
      const list = el('div', { class: 'wz-filelist' });
      state[side].files.forEach((f, idx) => {
        const ext = (f.name.split('.').pop() || '').toLowerCase();
        const row = el('div', { class: 'wz-file' }, [
          el('span', { class: 'wz-ext' }, ext || 'bin'),
          el('span', { class: 'wz-name' }, f.name + ' · ' + humanSize(f.size)),
          el('span', { class: 'wz-status ' + (f.err ? 'err' : 'ok') }, f.err ? 'Error: ' + f.err : (f.tables ? f.tables.length + ' table(s)' : 'parsed')),
          el('button', { class: 'wz-rem', onclick: () => { state[side].files.splice(idx, 1); recomputeTables(side); renderBody(); } }, '×'),
        ]);
        list.appendChild(row);
      });
      wrap.appendChild(list);
    }
    return wrap;
  }

  function humanSize(b) {
    if (b < 1024) return b + ' B';
    if (b < 1024*1024) return (b/1024).toFixed(1) + ' KB';
    return (b/1024/1024).toFixed(2) + ' MB';
  }

  async function handleFiles(side, fileList) {
    const files = Array.from(fileList || []);
    if (!files.length) return;
    toast('Parsing ' + files.length + ' file(s)…');
    for (const f of files) {
      const entry = { name: f.name, size: f.size, tables: null, err: null, sha256: null };
      state[side].files.push(entry);
      renderBody();
      try {
        const buf = await f.arrayBuffer();
        entry.sha256 = await sha256Hex(buf);
        const tables = await parseFile(f, buf);
        if (!tables || !tables.length) throw new Error('No tabular data found in file');
        entry.tables = tables;
      } catch (e) {
        entry.err = e.message || String(e);
        toast('Failed to parse ' + f.name + ': ' + entry.err, 'err');
      }
    }
    recomputeTables(side);
    renderBody();
  }

  function recomputeTables(side) {
    const all = [];
    state[side].files.forEach(f => { if (f.tables) f.tables.forEach(t => all.push({ fileName: f.name, sha256: f.sha256, ...t })); });
    state[side].tables = all;
    state[side].selectedTable = 0;
    state[side].records = all.length ? all[0].rows : [];
    state[side].status = all.length ? 'parsed' : 'empty';
    // Aggregate sha256 for multiple files
    state[side].sha256 = state[side].files.map(f => f.sha256).filter(Boolean).join(',');
  }

  async function parseFile(file, buf) {
    const ext = (file.name.split('.').pop() || '').toLowerCase();
    if (ext === 'csv' || ext === 'tsv' || ext === 'txt') return parseCSVBuf(buf, ext === 'tsv' ? '\t' : (ext === 'txt' ? null : ','));
    if (ext === 'json') return parseJSONBuf(buf);
    if (ext === 'xlsx' || ext === 'xls') {
      if (!(await ensureSheetJS())) throw new Error('SheetJS library unavailable (CDN blocked). Try CSV instead, or attach the file in chat for server-side parsing.');
      return parseXLSXBuf(buf);
    }
    if (ext === 'docx') {
      if (!(await ensureMammoth())) throw new Error('mammoth.js unavailable (CDN blocked). Export the tables to CSV, or attach the file in chat.');
      return parseDOCXBuf(buf);
    }
    if (ext === 'pptx') {
      if (!(await ensureJSZip())) throw new Error('JSZip unavailable (CDN blocked). Export the tables to CSV, or attach the file in chat.');
      return parsePPTXBuf(buf);
    }
    if (ext === 'pdf') {
      if (!(await ensurePdfJs())) throw new Error('pdf.js unavailable (CDN blocked). Export the tables to CSV, or attach the file in chat.');
      return parsePDFBuf(buf);
    }
    throw new Error('Unsupported file extension: ' + ext);
  }

  function decodeBuf(buf) {
    try { return new TextDecoder('utf-8', { fatal: false }).decode(buf); }
    catch (e) { return new TextDecoder('latin1').decode(buf); }
  }

  function parseCSVBuf(buf, sep) {
    const text = decodeBuf(buf);
    if (sep == null) sep = (text.split('\n')[0].match(/\t/) ? '\t' : ',');
    const rows = parseCSVText(text, sep);
    if (rows.length < 2) return [];
    const headers = rows[0];
    const data = rows.slice(1).filter(r => r.length && r.some(c => c !== ''))
      .map(r => { const o = {}; headers.forEach((h, i) => o[h || ('col' + (i+1))] = r[i] == null ? '' : r[i]); return o; });
    return [{ name: 'data', rows: data }];
  }

  function parseCSVText(text, sep) {
    const out = [];
    let row = [], field = '', i = 0, inQ = false;
    while (i < text.length) {
      const c = text[i];
      if (inQ) {
        if (c === '"' && text[i+1] === '"') { field += '"'; i += 2; continue; }
        if (c === '"') { inQ = false; i++; continue; }
        field += c; i++; continue;
      }
      if (c === '"') { inQ = true; i++; continue; }
      if (c === sep) { row.push(field); field = ''; i++; continue; }
      if (c === '\n' || c === '\r') {
        row.push(field); field = '';
        if (row.length) out.push(row);
        row = [];
        if (c === '\r' && text[i+1] === '\n') i++;
        i++; continue;
      }
      field += c; i++;
    }
    if (field.length || row.length) { row.push(field); out.push(row); }
    return out;
  }

  function parseJSONBuf(buf) {
    const text = decodeBuf(buf);
    const obj = JSON.parse(text);
    if (Array.isArray(obj)) return [{ name: 'records', rows: obj }];
    if (obj && Array.isArray(obj.records)) return [{ name: 'records', rows: obj.records }];
    if (obj && typeof obj === 'object') {
      const out = [];
      for (const k in obj) if (Array.isArray(obj[k])) out.push({ name: k, rows: obj[k] });
      if (out.length) return out;
    }
    throw new Error('JSON must be an array or an object with array properties.');
  }

  function parseXLSXBuf(buf) {
    const wb = XLSX.read(buf, { type: 'array' });
    const out = [];
    wb.SheetNames.forEach(name => {
      const ws = wb.Sheets[name];
      const rows = XLSX.utils.sheet_to_json(ws, { defval: '', raw: false });
      if (rows.length) out.push({ name: name, rows: rows });
    });
    return out;
  }

  async function parseDOCXBuf(buf) {
    const result = await mammoth.convertToHtml({ arrayBuffer: buf });
    const parser = new DOMParser();
    const doc = parser.parseFromString(result.value, 'text/html');
    const tables = Array.from(doc.querySelectorAll('table'));
    return tables.map((t, i) => ({ name: 'table_' + (i + 1), rows: tableEltToRecords(t) })).filter(x => x.rows.length);
  }

  function tableEltToRecords(tbl) {
    const trs = Array.from(tbl.querySelectorAll('tr'));
    if (!trs.length) return [];
    const cellTxt = r => Array.from(r.querySelectorAll('th,td')).map(c => (c.textContent || '').trim());
    const headers = cellTxt(trs[0]);
    return trs.slice(1).map(r => {
      const c = cellTxt(r);
      const o = {};
      headers.forEach((h, i) => o[h || ('col' + (i+1))] = c[i] == null ? '' : c[i]);
      return o;
    });
  }

  async function parsePPTXBuf(buf) {
    const zip = await JSZip.loadAsync(buf);
    const slideNames = Object.keys(zip.files).filter(n => /^ppt\/slides\/slide\d+\.xml$/i.test(n)).sort();
    const out = [];
    for (const n of slideNames) {
      const xml = await zip.files[n].async('string');
      const dom = new DOMParser().parseFromString(xml, 'application/xml');
      const tables = Array.from(dom.getElementsByTagName('a:tbl'));
      tables.forEach((tbl, idx) => {
        const rows = Array.from(tbl.getElementsByTagName('a:tr')).map(tr =>
          Array.from(tr.getElementsByTagName('a:tc')).map(tc =>
            Array.from(tc.getElementsByTagName('a:t')).map(t => t.textContent || '').join(' ').trim()
          )
        );
        if (rows.length > 1) {
          const headers = rows[0];
          const records = rows.slice(1).map(r => { const o = {}; headers.forEach((h, i) => o[h || ('col' + (i+1))] = r[i] || ''); return o; });
          out.push({ name: n.replace(/^.*\//, '').replace('.xml', '') + (tables.length > 1 ? '_' + (idx + 1) : ''), rows: records });
        }
      });
    }
    return out;
  }

  async function parsePDFBuf(buf) {
    const pdf = await pdfjsLib.getDocument({ data: buf }).promise;
    const allRows = [];
    for (let p = 1; p <= pdf.numPages; p++) {
      const page = await pdf.getPage(p);
      const content = await page.getTextContent();
      const byY = new Map();
      content.items.forEach(it => {
        const y = Math.round(it.transform[5]);
        if (!byY.has(y)) byY.set(y, []);
        byY.get(y).push({ x: it.transform[4], text: (it.str || '').trim() });
      });
      const ys = Array.from(byY.keys()).sort((a, b) => b - a);
      ys.forEach(y => {
        const cells = byY.get(y).sort((a, b) => a.x - b.x).map(c => c.text).filter(s => s);
        if (cells.length > 1) allRows.push(cells);
      });
    }
    if (allRows.length < 2) return [];
    // Heuristic: first row with most columns wins as header
    let headerIdx = 0, maxCols = 0;
    allRows.forEach((r, i) => { if (r.length > maxCols) { maxCols = r.length; headerIdx = i; } });
    const headers = allRows[headerIdx];
    const rows = allRows.slice(headerIdx + 1)
      .filter(r => r.length >= Math.max(2, Math.floor(maxCols / 2)))
      .map(r => { const o = {}; headers.forEach((h, i) => o[h || ('col' + (i+1))] = r[i] == null ? '' : r[i]); return o; });
    return rows.length ? [{ name: 'pdf_table', rows: rows, pdfNote: true }] : [];
  }

  function renderTablePreview(side) {
    const s = state[side];
    const box = el('div', { class: 'wz-preview' });
    const header = el('header', {});
    header.appendChild(el('span', {}, (side === 'left' ? 'Left' : 'Right') + ' — ' + s.records.length + ' records detected'));
    if (s.tables.length > 1) {
      const sel = el('div', { class: 'wz-tabselect' });
      s.tables.forEach((t, i) => {
        const tab = el('span', { class: 'wz-tab' + (i === s.selectedTable ? ' active' : ''), onclick: () => { s.selectedTable = i; s.records = s.tables[i].rows; renderBody(); } }, (t.fileName ? t.fileName + ' · ' : '') + t.name);
        sel.appendChild(tab);
      });
      header.appendChild(sel);
    }
    box.appendChild(header);
    const scrollDiv = el('div', { class: 'wz-scroll' });
    const tbl = buildPreviewTable(s.records.slice(0, 8));
    scrollDiv.appendChild(tbl);
    box.appendChild(scrollDiv);
    if (s.tables[s.selectedTable] && s.tables[s.selectedTable].pdfNote) {
      box.appendChild(el('div', { class: 'wz-banner warn', style: 'margin:10px 14px;', html:
        '<div>PDF table extraction is best-effort. If columns look misaligned, re-upload this side as CSV or XLSX for highest fidelity.</div>'
      }));
    }
    return box;
  }

  function buildPreviewTable(rows) {
    const tbl = el('table', {});
    if (!rows.length) { tbl.appendChild(el('tbody', {}, el('tr', {}, el('td', {}, '(empty)')))); return tbl; }
    const keys = Object.keys(rows[0]);
    const thead = el('thead', {});
    const trh = el('tr', {});
    keys.forEach(k => trh.appendChild(el('th', {}, k)));
    thead.appendChild(trh); tbl.appendChild(thead);
    const tbody = el('tbody', {});
    rows.forEach(r => {
      const tr = el('tr', {});
      keys.forEach(k => tr.appendChild(el('td', {}, String(r[k] == null ? '' : r[k]))));
      tbody.appendChild(tr);
    });
    tbl.appendChild(tbody);
    return tbl;
  }

  // --- Step 3: Rules ---------------------------------------------------------
  function renderStep3() {
    const panel = el('div', { class: 'wz-panel' });
    panel.appendChild(el('h2', {}, STEPS[2].tag + ' — ' + STEPS[2].label));
    panel.appendChild(el('p', { class: 'wz-sub' }, 'We auto-detected candidate key and amount fields from your data. Refine them or accept defaults.'));

    const lkeys = state.left.records.length ? Object.keys(state.left.records[0]) : [];
    const rkeys = state.right.records.length ? Object.keys(state.right.records[0]) : [];
    const commonKeys = lkeys.filter(k => rkeys.includes(k));

    if (!state.rules.keyFields.length) {
      const defaults = SECTORS[state.scope.sector] || SECTORS['Other'];
      state.rules.keyFields = (defaults.key_fields || []).filter(k => commonKeys.includes(k));
      if (!state.rules.keyFields.length) state.rules.keyFields = commonKeys.filter(isIdKey).slice(0, 2);
      if (!state.rules.keyFields.length && commonKeys.length) state.rules.keyFields = [commonKeys[0]];
    }
    if (!state.rules.amountFields.length) {
      const defaults = SECTORS[state.scope.sector] || SECTORS['Other'];
      state.rules.amountFields = (defaults.amount_fields || []).filter(k => commonKeys.includes(k));
      if (!state.rules.amountFields.length) state.rules.amountFields = commonKeys.filter(isAmountKey).slice(0, 3);
    }

    const grid = el('div', { class: 'wz-grid' });
    grid.appendChild(field('Key fields (match identifier)',
      multiTagSelect(commonKeys, state.rules.keyFields, list => state.rules.keyFields = list),
      'Columns used to join left and right. Must exist on both sides.'
    ));
    grid.appendChild(field('Amount fields (numeric comparison)',
      multiTagSelect(commonKeys, state.rules.amountFields, list => state.rules.amountFields = list),
      'Numeric columns checked for variance using the tolerance below.'
    ));
    grid.appendChild(field('Absolute tolerance',
      inputNode('number', String(state.rules.tolerance), v => state.rules.tolerance = parseFloat(v) || 0),
      'Two amounts within this absolute difference are treated as a match.'
    ));
    grid.appendChild(field('Percent tolerance (optional)',
      inputNode('number', String(state.rules.tolerancePct), v => state.rules.tolerancePct = parseFloat(v) || 0),
      'Percent-of-nominal tolerance applied in addition to absolute.'
    ));
    grid.appendChild(field('Currency mode',
      selectInput(['strict', 'ignore'], state.rules.currencyMode, v => state.rules.currencyMode = v),
      '"strict" never nets across currencies; "ignore" compares numerics only.'
    ));
    grid.appendChild(field('PII redaction (HIPAA / GDPR)',
      toggleInput(state.rules.piiRedact, v => state.rules.piiRedact = v, 'Mask names, DOB, and IDs in downloads'),
      'Auto-enabled for Healthcare and Pharma/Clinical sectors.'
    ));
    panel.appendChild(grid);

    const unmatchedLeft = state.rules.keyFields.filter(k => !lkeys.includes(k));
    const unmatchedRight = state.rules.keyFields.filter(k => !rkeys.includes(k));
    if (unmatchedLeft.length || unmatchedRight.length) {
      panel.appendChild(el('div', { class: 'wz-banner err', html:
        '<div>Key field(s) missing: ' +
        (unmatchedLeft.length ? '<b>Left</b> is missing ' + esc(unmatchedLeft.join(', ')) + '. ' : '') +
        (unmatchedRight.length ? '<b>Right</b> is missing ' + esc(unmatchedRight.join(', ')) + '.' : '') +
        ' Pick columns present in both datasets.</div>'
      }));
    }

    const canNext = state.rules.keyFields.length > 0 && !unmatchedLeft.length && !unmatchedRight.length;
    panel.appendChild(renderActions(true, canNext, 'Run reconciliation'));
    return panel;
  }

  function multiTagSelect(options, value, onchange) {
    const wrap = el('div', {});
    const row = el('div', { class: 'wz-tagrow' });
    value.forEach((v, idx) => {
      const tag = el('span', { class: 'wz-tag removable' }, [
        v,
        el('button', { class: 'wz-x', onclick: () => { value.splice(idx, 1); onchange(value); renderBody(); } }, '×'),
      ]);
      row.appendChild(tag);
    });
    wrap.appendChild(row);
    const sel = el('select', { class: 'wz-select', style: 'margin-top:6px;' });
    sel.appendChild(el('option', { value: '' }, '+ Add column…'));
    options.filter(o => !value.includes(o)).forEach(o => sel.appendChild(el('option', { value: o }, o)));
    sel.addEventListener('change', () => {
      if (sel.value) { value.push(sel.value); onchange(value); renderBody(); }
    });
    wrap.appendChild(sel);
    return wrap;
  }

  function toggleInput(checked, onchange, label) {
    const w = el('label', { class: 'wz-switch', style: 'margin-top:8px;' });
    const inp = el('input', { type: 'checkbox' }); inp.checked = !!checked;
    inp.addEventListener('change', () => onchange(inp.checked));
    w.appendChild(inp);
    w.appendChild(el('span', { class: 'wz-track' }));
    w.appendChild(el('span', { class: 'wz-switch-lbl' }, label || (checked ? 'On' : 'Off')));
    return w;
  }

  // --- Step 4: Run -----------------------------------------------------------
  function renderStep4() {
    const panel = el('div', { class: 'wz-panel' });
    panel.appendChild(el('h2', {}, STEPS[3].tag + ' — ' + STEPS[3].label));
    panel.appendChild(el('p', { class: 'wz-sub' }, 'Deterministic match on ' + esc(state.rules.keyFields.join(' + ')) + ' with variance check on ' + (state.rules.amountFields.length ? esc(state.rules.amountFields.join(', ')) : 'no amount fields') + '.'));

    const progWrap = el('div', {});
    const phases = ['Normalising', 'Indexing right side', 'Matching left→right', 'Variance check', 'Redaction & summary'];
    const rows = phases.map(() => {
      const row = el('div', { class: 'wz-progress-row' });
      row.appendChild(el('span', { class: 'wz-lbl' }, ''));
      const bar = el('div', { class: 'wz-progress' }, el('div', { class: 'wz-progress-bar' }));
      row.appendChild(bar);
      row.appendChild(el('span', { class: 'wz-val' }, '0%'));
      progWrap.appendChild(row);
      return row;
    });
    rows.forEach((r, i) => $('.wz-lbl', r).textContent = phases[i]);
    panel.appendChild(progWrap);

    const runBtn = el('button', { class: 'wz-btn', style: 'margin-top:16px;', onclick: () => runReconcile(rows) }, 'Start reconciliation');
    panel.appendChild(el('div', { style: 'margin-top:14px;' }, runBtn));

    panel.appendChild(renderActions(true, !!state.result, state.result ? 'Next — review & download' : 'Run to continue'));
    return panel;
  }

  async function runReconcile(rows) {
    const setPhase = (i, pct) => {
      const bar = $('.wz-progress-bar', rows[i]);
      const val = $('.wz-val', rows[i]);
      bar.style.width = pct + '%';
      val.textContent = pct + '%';
    };
    try {
      for (let i = 0; i < rows.length; i++) setPhase(i, 0);
      setPhase(0, 20);
      const left = state.left.records.slice();
      const right = state.right.records.slice();
      await sleep(40);
      setPhase(0, 100);

      setPhase(1, 30);
      const kf = state.rules.keyFields;
      const af = state.rules.amountFields;
      const tol = state.rules.tolerance;
      const tolPct = state.rules.tolerancePct;
      const rightIdx = new Map();
      right.forEach((r, idx) => {
        const k = keyOf(r, kf);
        if (!rightIdx.has(k)) rightIdx.set(k, []);
        rightIdx.get(k).push(idx);
      });
      setPhase(1, 100);

      setPhase(2, 10);
      const used = new Set();
      const matched = [], variance = [], unmatchedLeft = [];
      for (let i = 0; i < left.length; i++) {
        const lrec = left[i];
        const k = keyOf(lrec, kf);
        const cands = rightIdx.get(k) || [];
        let picked = -1;
        for (const ri of cands) { if (!used.has(ri)) { picked = ri; break; } }
        if (picked < 0) { unmatchedLeft.push(lrec); continue; }
        used.add(picked);
        const rrec = right[picked];
        const diffs = {};
        for (const f of af) {
          const lv = toNum(lrec[f]);
          const rv = toNum(rrec[f]);
          if (lv == null && rv == null) continue;
          if (lv == null || rv == null) { diffs[f] = { left: lv, right: rv, delta: null }; continue; }
          const abs = Math.abs(lv - rv);
          const pctOk = tolPct > 0 && Math.max(Math.abs(lv), Math.abs(rv)) > 0 && (abs / Math.max(Math.abs(lv), Math.abs(rv)) * 100) <= tolPct;
          if (abs > tol && !pctOk) diffs[f] = { left: lv, right: rv, delta: Math.round((rv - lv) * 10000) / 10000 };
        }
        if (Object.keys(diffs).length) {
          const row = { _key: k };
          for (const kk in lrec) row[kk] = lrec[kk];
          for (const ff in diffs) { row[ff + '_left'] = diffs[ff].left; row[ff + '_right'] = diffs[ff].right; row[ff + '_delta'] = diffs[ff].delta; }
          variance.push(row);
        } else {
          matched.push(lrec);
        }
        if (i % 200 === 0) { setPhase(2, Math.min(99, 10 + Math.round((i / Math.max(1, left.length)) * 85))); await sleep(0); }
      }
      setPhase(2, 100);

      setPhase(3, 50);
      const unmatchedRight = right.filter((_, idx) => !used.has(idx));
      setPhase(3, 100);

      setPhase(4, 40);
      let redLeft = unmatchedLeft, redRight = unmatchedRight, redVar = variance, redMat = matched;
      if (state.rules.piiRedact) {
        redLeft = redactPII(unmatchedLeft);
        redRight = redactPII(unmatchedRight);
        redVar = redactPII(variance);
        redMat = redactPII(matched);
      }
      const total = matched.length + variance.length + unmatchedLeft.length;
      const match_rate = total ? Math.round((matched.length / total) * 10000) / 100 : 0;
      state.result = {
        matched: redMat, variance: redVar,
        unmatched_left: redLeft, unmatched_right: redRight,
        stats: {
          total_left: left.length, total_right: right.length,
          matched: matched.length, variance: variance.length,
          unmatched_left: unmatchedLeft.length, unmatched_right: unmatchedRight.length,
          match_rate: match_rate,
        },
        ranAt: new Date().toISOString(),
      };
      setPhase(4, 100);
      await sleep(120);
      toast('Reconciliation complete — ' + match_rate + '% match rate', 'ok');
      goto(5);
    } catch (e) {
      console.error(e);
      toast('Reconciliation failed: ' + (e.message || String(e)), 'err');
    }
  }

  function sleep(ms) { return new Promise(r => setTimeout(r, ms)); }
  function keyOf(rec, kf) { return kf.map(k => normStr(rec[k])).join('||'); }
  function normStr(v) { return String(v == null ? '' : v).replace(/\s+/g, ' ').trim().toLowerCase(); }
  function toNum(v) {
    if (v == null || v === '') return null;
    const n = parseFloat(String(v).replace(/,/g, '').replace(/[^0-9.\-]/g, ''));
    return isFinite(n) ? n : null;
  }

  function redactPII(rows) {
    const nameKeys = /(patient.?name|insured|policyholder|subject.?name|patientfullname|full.?name)/i;
    const dobKeys = /(patient.?dob|dob|birth.?date|date.?of.?birth)/i;
    const idKeys = /(ssn|social|member.?id|patient.?id|mrn|national.?id)/i;
    return rows.map(r => {
      const o = {};
      for (const k in r) {
        const v = r[k];
        if (nameKeys.test(k) && v) {
          const parts = String(v).trim().split(/[\s,]+/);
          o[k] = parts.map(p => p ? p[0] + '.' : '').join(' ') || '***';
        } else if (dobKeys.test(k) && v) {
          const s = String(v); o[k] = s.length >= 4 ? s.slice(0, 4) + '-**-**' : '****';
        } else if (idKeys.test(k) && v) {
          const s = String(v); o[k] = s.length > 5 ? s.slice(0, 3) + '***' + s.slice(-2) : '***';
        } else {
          o[k] = v;
        }
      }
      return o;
    });
  }

  // --- Step 5: Review & Download --------------------------------------------
  function renderStep5() {
    const panel = el('div', { class: 'wz-panel' });
    panel.appendChild(el('h2', {}, STEPS[4].tag + ' — ' + STEPS[4].label));
    if (!state.result) {
      panel.appendChild(el('div', { class: 'wz-banner warn' }, 'Reconciliation has not been run yet. Go back to Step 4.'));
      panel.appendChild(renderActions(true, false, 'Run first'));
      return panel;
    }
    const r = state.result, s = r.stats;
    panel.appendChild(el('p', { class: 'wz-sub' }, 'Reconciliation complete. Downloads are generated in your browser; outputs include the run ID, SHA-256 of both inputs, operator name, timestamp, and regulation tags for audit traceability.'));

    // Regulation tags
    const regs = ((SECTORS[state.scope.sector] || SECTORS['Other'] || {}).regulations || {})[state.scope.region] || [];
    if (regs.length) {
      const tagRow = el('div', { class: 'wz-tagrow', style: 'margin-bottom:12px;' });
      regs.forEach(t => tagRow.appendChild(el('span', { class: 'wz-tag' }, t)));
      panel.appendChild(tagRow);
    }

    // KPIs
    const kpis = el('div', { class: 'wz-kpis' });
    const kpi = (v, l, cls) => { const k = el('div', { class: 'wz-kpi ' + (cls || '') }); k.appendChild(el('div', { class: 'v' }, fmtNum(v))); k.appendChild(el('div', { class: 'l' }, l)); return k; };
    kpis.appendChild(kpi(s.match_rate + '%', 'Match rate', s.match_rate >= 95 ? 'ok' : (s.match_rate >= 80 ? 'warn' : 'err')));
    kpis.appendChild(kpi(s.total_left, 'Records · Left'));
    kpis.appendChild(kpi(s.total_right, 'Records · Right'));
    kpis.appendChild(kpi(s.matched, 'Matched', 'ok'));
    kpis.appendChild(kpi(s.variance, 'Variance', 'warn'));
    kpis.appendChild(kpi(s.unmatched_left, 'Missing · Right', s.unmatched_left > 0 ? 'err' : ''));
    kpis.appendChild(kpi(s.unmatched_right, 'Missing · Left', s.unmatched_right > 0 ? 'err' : ''));
    panel.appendChild(kpis);

    // Download bar
    const dl = el('div', { class: 'wz-dl-bar' });
    dl.appendChild(el('span', { class: 'wz-lbl' }, 'Download reconciled output:'));
    dl.appendChild(el('button', { class: 'wz-btn', onclick: () => downloadXLSX() }, 'XLSX'));
    dl.appendChild(el('button', { class: 'wz-btn wz-secondary', onclick: () => downloadCSV() }, 'CSV (exceptions)'));
    dl.appendChild(el('button', { class: 'wz-btn wz-secondary', onclick: () => downloadPDF() }, 'PDF (audit memo)'));
    dl.appendChild(el('button', { class: 'wz-btn wz-secondary', onclick: () => downloadDOCX() }, 'DOCX (narrative)'));
    dl.appendChild(el('button', { class: 'wz-btn wz-ghost', onclick: () => downloadJSON() }, 'JSON'));
    panel.appendChild(dl);

    // Tables
    if (r.variance.length) panel.appendChild(renderResultTable('Variances — amount mismatches', r.variance));
    if (r.unmatched_left.length) panel.appendChild(renderResultTable('Missing in Right (left-only)', r.unmatched_left));
    if (r.unmatched_right.length) panel.appendChild(renderResultTable('Missing in Left (right-only)', r.unmatched_right));
    if (r.matched.length) panel.appendChild(renderResultTable('Matched (first 25)', r.matched.slice(0, 25)));

    // Sign-off block
    const sign = el('div', { class: 'wz-signoff' });
    sign.innerHTML = [
      '<b>Run ID:</b> ' + esc(state.scope.runId),
      '<b>Timestamp:</b> ' + esc(r.ranAt),
      '<b>Operator:</b> ' + esc(state.scope.operator || '—'),
      '<b>Sector / Region:</b> ' + esc(state.scope.sector) + ' · ' + esc(state.scope.region),
      '<b>As of:</b> ' + esc(state.scope.asOf),
      '<b>Key fields:</b> ' + esc(state.rules.keyFields.join(' + ')),
      '<b>Amount fields:</b> ' + esc(state.rules.amountFields.join(', ') || '—'),
      '<b>Tolerance:</b> abs ' + state.rules.tolerance + (state.rules.tolerancePct ? ' · pct ' + state.rules.tolerancePct + '%' : ''),
      '<b>Left SHA-256:</b> ' + esc(state.left.sha256 || '—'),
      '<b>Right SHA-256:</b> ' + esc(state.right.sha256 || '—'),
      '<b>Regulation tags:</b> ' + esc(regs.join(', ') || '—'),
      '<b>PII redaction:</b> ' + (state.rules.piiRedact ? 'applied' : 'not applied'),
      '<b>Build:</b> ' + esc(BUILD),
    ].join('<br>');
    panel.appendChild(sign);

    panel.appendChild(renderActions(true, false, null, true));
    return panel;
  }

  function renderResultTable(title, rows) {
    const box = el('div', { class: 'wz-tbl-wrap' });
    const h = el('header', {});
    h.appendChild(el('span', {}, title));
    h.appendChild(el('span', { class: 'wz-count' }, rows.length + ' row' + (rows.length === 1 ? '' : 's')));
    box.appendChild(h);
    const scrollDiv = el('div', { class: 'wz-scroll' });
    if (!rows.length) {
      scrollDiv.appendChild(el('div', { style: 'padding:14px;color:var(--wz-muted);' }, 'None.'));
    } else {
      const keys = Object.keys(rows[0]);
      const tbl = el('table', {});
      const trh = el('tr', {});
      keys.forEach(k => trh.appendChild(el('th', { class: isAmountKey(k) ? 'num' : '' }, k)));
      tbl.appendChild(el('thead', {}, trh));
      const tb = el('tbody', {});
      rows.slice(0, 200).forEach(row => {
        const tr = el('tr', {});
        keys.forEach(k => {
          const v = row[k];
          let cls = isAmountKey(k) ? 'num' : '';
          if (/_delta$/.test(k)) { const n = parseFloat(v); if (isFinite(n)) cls += (n < 0 ? ' delta-neg' : (n > 0 ? ' delta-pos' : '')); }
          tr.appendChild(el('td', { class: cls }, isAmountKey(k) ? fmtNum(v) : (v == null ? '' : String(v))));
        });
        tb.appendChild(tr);
      });
      tbl.appendChild(tb);
      scrollDiv.appendChild(tbl);
    }
    box.appendChild(scrollDiv);
    return box;
  }

  // --- Actions bar -----------------------------------------------------------
  function renderActions(showBack, canNext, nextLabel, isFinal) {
    const bar = el('div', { class: 'wz-actions' });
    const left = el('div', {});
    if (showBack) left.appendChild(el('button', { class: 'wz-btn wz-ghost', onclick: () => goto(state.step - 1) }, '← Back'));
    bar.appendChild(left);
    const right = el('div', { style: 'display:flex;gap:10px;align-items:center;' });
    if (isFinal) {
      right.appendChild(el('button', { class: 'wz-btn wz-secondary', onclick: () => { if (confirm('Start a new reconciliation? Current results will be cleared.')) resetAll(); } }, 'Start new run'));
    } else if (nextLabel) {
      const btn = el('button', { class: 'wz-btn', onclick: () => { if (canNext) goto(state.step + 1); } }, nextLabel);
      if (!canNext) btn.disabled = true;
      right.appendChild(btn);
    }
    bar.appendChild(right);
    return bar;
  }

  function resetAll() {
    state.step = 1;
    state.left = { files: [], tables: [], selectedTable: 0, records: [], sha256: null, status: 'empty' };
    state.right = { files: [], tables: [], selectedTable: 0, records: [], sha256: null, status: 'empty' };
    state.result = null;
    state.scope.runId = genRunId();
    state.rules.keyFields = []; state.rules.amountFields = [];
    renderStepper(); renderBody();
  }

  // --- Download generators ---------------------------------------------------
  function baseName() { return (state.scope.outputName || 'reconciliation') + '_' + state.scope.asOf.replace(/-/g, ''); }

  function downloadBlob(name, mime, data) {
    const blob = data instanceof Blob ? data : new Blob([data], { type: mime });
    const url = URL.createObjectURL(blob);
    const a = el('a', { href: url, download: name });
    document.body.appendChild(a); a.click();
    setTimeout(() => { URL.revokeObjectURL(url); a.remove(); }, 400);
  }

  function downloadCSV() {
    try {
      const parts = [];
      const add = (title, rows) => {
        parts.push('# ' + title);
        if (!rows.length) { parts.push(''); return; }
        const keys = Object.keys(rows[0]);
        parts.push(keys.map(csvCell).join(','));
        rows.forEach(r => parts.push(keys.map(k => csvCell(r[k])).join(',')));
        parts.push('');
      };
      const r = state.result;
      add('Summary', [Object.assign({ run_id: state.scope.runId, operator: state.scope.operator, asOf: state.scope.asOf, sector: state.scope.sector, region: state.scope.region }, r.stats)]);
      add('Variance', r.variance);
      add('Missing_in_Right', r.unmatched_left);
      add('Missing_in_Left', r.unmatched_right);
      add('Matched (first 500)', r.matched.slice(0, 500));
      downloadBlob(baseName() + '.csv', 'text/csv', parts.join('\n'));
      toast('CSV downloaded', 'ok');
    } catch (e) { toast('CSV failed: ' + e.message, 'err'); }
  }
  function csvCell(v) {
    const s = v == null ? '' : String(v);
    return /[",\n\r]/.test(s) ? '"' + s.replace(/"/g, '""') + '"' : s;
  }

  function downloadJSON() {
    try {
      const payload = {
        meta: { runId: state.scope.runId, operator: state.scope.operator, asOf: state.scope.asOf, sector: state.scope.sector, region: state.scope.region, keyFields: state.rules.keyFields, amountFields: state.rules.amountFields, tolerance: state.rules.tolerance, leftSHA256: state.left.sha256, rightSHA256: state.right.sha256, build: BUILD },
        result: state.result,
      };
      downloadBlob(baseName() + '.json', 'application/json', JSON.stringify(payload, null, 2));
      toast('JSON downloaded', 'ok');
    } catch (e) { toast('JSON failed: ' + e.message, 'err'); }
  }

  async function downloadXLSX() {
    if (!(await ensureSheetJS())) { toast('XLSX library unavailable — download CSV instead, or attach files in chat for server-side ooXML', 'err'); return; }
    try {
      const wb = XLSX.utils.book_new();
      const regs = ((SECTORS[state.scope.sector] || SECTORS['Other'] || {}).regulations || {})[state.scope.region] || [];
      const summary = [
        ['IBM Consulting Advantage — AI Reconciliation'],
        [],
        ['Sector', state.scope.sector],
        ['Region', state.scope.region],
        ['As of', state.scope.asOf],
        ['Run ID', state.scope.runId],
        ['Operator', state.scope.operator || '—'],
        ['Timestamp', state.result.ranAt],
        ['Left SHA-256', state.left.sha256 || '—'],
        ['Right SHA-256', state.right.sha256 || '—'],
        ['Key fields', state.rules.keyFields.join(' + ')],
        ['Amount fields', state.rules.amountFields.join(', ')],
        ['Absolute tolerance', state.rules.tolerance],
        ['Percent tolerance', state.rules.tolerancePct],
        ['PII redaction', state.rules.piiRedact ? 'applied' : 'not applied'],
        ['Regulation tags', regs.join(', ')],
        ['Build', BUILD],
        [],
        ['KPI', 'Value'],
      ];
      const s = state.result.stats;
      Object.keys(s).forEach(k => summary.push([k, s[k]]));
      XLSX.utils.book_append_sheet(wb, XLSX.utils.aoa_to_sheet(summary), 'Summary');
      const secs = [['Variance', state.result.variance], ['Missing_in_Right', state.result.unmatched_left], ['Missing_in_Left', state.result.unmatched_right], ['Matched', state.result.matched]];
      secs.forEach(([name, rows]) => {
        if (!rows.length) return;
        XLSX.utils.book_append_sheet(wb, XLSX.utils.json_to_sheet(rows), name.slice(0, 31));
      });
      XLSX.writeFile(wb, baseName() + '.xlsx');
      toast('XLSX downloaded', 'ok');
    } catch (e) { toast('XLSX failed: ' + e.message, 'err'); }
  }

  async function downloadPDF() {
    if (!(await ensureJsPDF())) { toast('PDF library unavailable — try XLSX or CSV, or attach files in chat', 'err'); return; }
    const jsPDFCtor = (window.jspdf || {}).jsPDF;
    try {
      const doc = new jsPDFCtor({ orientation: 'p', unit: 'pt', format: 'a4' });
      const NAVY = '#002D74';
      const regs = ((SECTORS[state.scope.sector] || SECTORS['Other'] || {}).regulations || {})[state.scope.region] || [];
      doc.setFillColor(NAVY); doc.rect(0, 0, doc.internal.pageSize.getWidth(), 70, 'F');
      doc.setTextColor(255); doc.setFont('helvetica', 'bold'); doc.setFontSize(16);
      doc.text('IBM Consulting Advantage — AI Reconciliation', 36, 30);
      doc.setFont('helvetica', 'normal'); doc.setFontSize(10);
      doc.text(state.scope.sector + ' · ' + state.scope.region + ' · As of ' + state.scope.asOf + ' · Run ' + state.scope.runId, 36, 50);

      let y = 92;
      doc.setTextColor(0); doc.setFont('helvetica', 'bold'); doc.setFontSize(12); doc.text('Summary', 36, y); y += 12;
      const s = state.result.stats;
      const sumRows = [
        ['Match rate', s.match_rate + '%'],
        ['Total · Left', String(s.total_left)],
        ['Total · Right', String(s.total_right)],
        ['Matched', String(s.matched)],
        ['Variance', String(s.variance)],
        ['Missing in Right', String(s.unmatched_left)],
        ['Missing in Left', String(s.unmatched_right)],
        ['Operator', state.scope.operator || '—'],
        ['Timestamp', state.result.ranAt],
        ['Regulation tags', regs.join(', ') || '—'],
      ];
      if (doc.autoTable) {
        doc.autoTable({ startY: y, head: [['Field', 'Value']], body: sumRows, theme: 'grid', styles: { fontSize: 9 }, headStyles: { fillColor: NAVY, textColor: 255 }, margin: { left: 36, right: 36 } });
        y = doc.lastAutoTable.finalY + 14;
      }

      const tableSection = (title, rows) => {
        if (!rows.length || !doc.autoTable) return;
        if (y > 700) { doc.addPage(); y = 60; }
        doc.setFont('helvetica', 'bold'); doc.setFontSize(11); doc.text(title + ' (' + rows.length + ')', 36, y); y += 8;
        const keys = Object.keys(rows[0]).slice(0, 6);
        const body = rows.slice(0, 60).map(r => keys.map(k => (r[k] == null ? '' : String(r[k])).slice(0, 40)));
        doc.autoTable({ startY: y, head: [keys], body: body, theme: 'striped', styles: { fontSize: 7.5, cellPadding: 3 }, headStyles: { fillColor: NAVY, textColor: 255 }, margin: { left: 36, right: 36 } });
        y = doc.lastAutoTable.finalY + 14;
      };
      tableSection('Variances', state.result.variance);
      tableSection('Missing in Right', state.result.unmatched_left);
      tableSection('Missing in Left', state.result.unmatched_right);

      // Sign-off
      if (y > 720) { doc.addPage(); y = 60; }
      doc.setFont('helvetica', 'bold'); doc.setFontSize(11); doc.text('Sign-off', 36, y); y += 10;
      doc.setFont('courier', 'normal'); doc.setFontSize(8);
      const sign = [
        'Run ID: ' + state.scope.runId,
        'Left SHA-256: ' + (state.left.sha256 || '—').slice(0, 64),
        'Right SHA-256: ' + (state.right.sha256 || '—').slice(0, 64),
        'Key fields: ' + state.rules.keyFields.join(' + '),
        'Tolerance (abs/pct): ' + state.rules.tolerance + ' / ' + state.rules.tolerancePct,
        'PII redaction: ' + (state.rules.piiRedact ? 'applied' : 'not applied'),
        'Build: ' + BUILD,
      ];
      sign.forEach(line => { doc.text(line, 36, y); y += 11; });
      doc.save(baseName() + '.pdf');
      toast('PDF downloaded', 'ok');
    } catch (e) { toast('PDF failed: ' + e.message, 'err'); }
  }

  async function downloadDOCX() {
    // Prefer docx.js; fallback to html-docx-js; both load on demand.
    if (await ensureDocxJs()) {
      try { buildDocxJs(window.docx); return; } catch (e) { console.warn('docx.js failed, falling back', e); }
    }
    if (await ensureHtmlDocx()) {
      try { buildDocxHtml(); return; } catch (e) { console.warn(e); }
    }
    toast('DOCX library unavailable — try XLSX or PDF, or attach files in chat', 'err');
  }

  function buildDocxJs(dx) {
    const { Document, Packer, Paragraph, TextRun, Table, TableRow, TableCell, HeadingLevel, AlignmentType, WidthType, BorderStyle } = dx;
    const NAVY = '002D74';
    const regs = ((SECTORS[state.scope.sector] || SECTORS['Other'] || {}).regulations || {})[state.scope.region] || [];
    const header = new Paragraph({ children: [new TextRun({ text: 'IBM Consulting Advantage — AI Reconciliation', bold: true, size: 36, color: NAVY })] });
    const sub = new Paragraph({ children: [new TextRun({ text: state.scope.sector + ' · ' + state.scope.region + ' · As of ' + state.scope.asOf + ' · Run ' + state.scope.runId, size: 20, color: '525252' })] });
    const s = state.result.stats;
    const kpis = new Paragraph({ children: [new TextRun({ text: 'Match rate: ' + s.match_rate + '% · Matched: ' + s.matched + ' · Variance: ' + s.variance + ' · Missing-R: ' + s.unmatched_left + ' · Missing-L: ' + s.unmatched_right, size: 20 })] });
    const regTags = new Paragraph({ children: [new TextRun({ text: 'Regulation tags: ' + (regs.join(', ') || '—'), size: 18, color: '525252' })] });

    const makeTable = (title, rows) => {
      if (!rows.length) return null;
      const keys = Object.keys(rows[0]).slice(0, 8);
      const mkCell = (text, head) => new TableCell({
        width: { size: 100 / keys.length * 100, type: WidthType.DXA },
        shading: head ? { fill: NAVY } : undefined,
        children: [new Paragraph({ children: [new TextRun({ text: String(text == null ? '' : text).slice(0, 120), bold: !!head, color: head ? 'FFFFFF' : '000000', size: 16 })] })],
      });
      const headRow = new TableRow({ children: keys.map(k => mkCell(k, true)) });
      const dataRows = rows.slice(0, 40).map(r => new TableRow({ children: keys.map(k => mkCell(r[k], false)) }));
      return [new Paragraph({ text: title + ' (' + rows.length + ')', heading: HeadingLevel.HEADING_2 }), new Table({ rows: [headRow, ...dataRows], width: { size: 100, type: WidthType.PERCENTAGE } })];
    };
    const sections = [header, sub, kpis, regTags];
    [['Variances', state.result.variance], ['Missing in Right', state.result.unmatched_left], ['Missing in Left', state.result.unmatched_right]].forEach(([t, rows]) => {
      const tbl = makeTable(t, rows); if (tbl) sections.push(...tbl);
    });
    sections.push(new Paragraph({ text: 'Sign-off', heading: HeadingLevel.HEADING_2 }));
    [
      'Run ID: ' + state.scope.runId,
      'Operator: ' + (state.scope.operator || '—'),
      'Timestamp: ' + state.result.ranAt,
      'Left SHA-256: ' + (state.left.sha256 || '—'),
      'Right SHA-256: ' + (state.right.sha256 || '—'),
      'Key fields: ' + state.rules.keyFields.join(' + '),
      'Tolerance: abs ' + state.rules.tolerance + ' / pct ' + state.rules.tolerancePct,
      'PII redaction: ' + (state.rules.piiRedact ? 'applied' : 'not applied'),
      'Build: ' + BUILD,
    ].forEach(line => sections.push(new Paragraph({ children: [new TextRun({ text: line, size: 16, font: 'Consolas', color: '3B5472' })] })));

    const doc = new Document({ sections: [{ children: sections }] });
    Packer.toBlob(doc).then(blob => { downloadBlob(baseName() + '.docx', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', blob); toast('DOCX downloaded', 'ok'); });
  }

  function buildDocxHtml() {
    const NAVY = '#002D74';
    const regs = ((SECTORS[state.scope.sector] || SECTORS['Other'] || {}).regulations || {})[state.scope.region] || [];
    const tbl = (title, rows) => {
      if (!rows.length) return '';
      const keys = Object.keys(rows[0]).slice(0, 8);
      const th = keys.map(k => '<th style="background:' + NAVY + ';color:#fff;padding:6px;border:1px solid #666;">' + esc(k) + '</th>').join('');
      const tr = rows.slice(0, 40).map(r => '<tr>' + keys.map(k => '<td style="padding:6px;border:1px solid #ccc;font-size:11px;">' + esc(String(r[k] == null ? '' : r[k])) + '</td>').join('') + '</tr>').join('');
      return '<h3 style="color:' + NAVY + ';">' + esc(title) + ' (' + rows.length + ')</h3><table style="border-collapse:collapse;width:100%;">' + th + tr + '</table>';
    };
    const html = '<!DOCTYPE html><html><body style="font-family:Segoe UI,Arial;color:#161616;">' +
      '<div style="background:' + NAVY + ';color:#fff;padding:16px;"><h1 style="margin:0;">IBM Consulting Advantage — AI Reconciliation</h1><div style="opacity:.85;font-size:12px;">' + esc(state.scope.sector + ' · ' + state.scope.region + ' · ' + state.scope.asOf) + '</div></div>' +
      '<p>Regulation tags: ' + esc(regs.join(', ') || '—') + '</p>' +
      tbl('Variances', state.result.variance) + tbl('Missing in Right', state.result.unmatched_left) + tbl('Missing in Left', state.result.unmatched_right) +
      '</body></html>';
    const blob = window.htmlDocx.asBlob(html);
    downloadBlob(baseName() + '.docx', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', blob);
    toast('DOCX downloaded (html-docx fallback)', 'ok');
  }

  // --- Height autosize -------------------------------------------------------
  function autoHeight() {
    try {
      const h = Math.max(document.documentElement.scrollHeight, document.body.scrollHeight);
      window.parent.postMessage({ type: 'iframe:resize', height: h }, '*');
    } catch (e) {}
  }

  // --- Boot ------------------------------------------------------------------
  function boot() {
    renderStepper();
    renderBody();
    setInterval(autoHeight, 700);
    window.addEventListener('load', autoHeight);
  }

  if (document.readyState === 'loading') document.addEventListener('DOMContentLoaded', boot);
  else boot();
})();
</script>
"""

# ---------------------------------------------------------------------------
# HTML shell
# ---------------------------------------------------------------------------


_WIZARD_SECTOR_ORDER = [
    "Banking",
    "Investment Banking",
    "Insurance",
    "Healthcare",
    "Asset Management",
    "Pharma Clinical",
    "Energy",
    "Telecommunications",
    "Retail",
    "Manufacturing",
    "Public Sector",
    "Technology",
    "Transportation",
    "Other",
]
_WIZARD_REGION_ORDER = ["USA", "European Union", "United Kingdom", "Global"]


def _wizard_config() -> Dict[str, Any]:
    sectors: Dict[str, Any] = {}
    for s in _WIZARD_SECTOR_ORDER:
        d = _sector_defaults(s)
        sectors[s] = {
            "key_fields": d["key_fields"],
            "amount_fields": d["amount_fields"],
            "regulations": d["regulations"],
        }
    return {
        "sectors": sectors,
        "regions": _WIZARD_REGION_ORDER,
        "build": _BUILD,
    }


def _wizard_shell_html() -> str:
    cfg_json = json.dumps(_wizard_config())
    return f"""<!DOCTYPE html>
<html lang="en" data-ibm-build="{_BUILD}">
<head>
<meta charset="utf-8">
<title>IBM Consulting Advantage — AI Reconciliation</title>
<link rel="preconnect" href="https://cdn.jsdelivr.net" crossorigin>
<style>{THEME_CSS}</style>
<style>{WIZARD_CSS}</style>
{CDN_SCRIPTS}
</head>
<body class="wz-root">
<div class="wz-shell">
  <div class="wz-head">
    <h1><span class="wz-dot"></span>IBM Consulting Advantage · AI Reconciliation</h1>
    <div class="wz-meta">Build {_BUILD}</div>
  </div>
  <div id="wz-stepper" class="wz-stepper"></div>
  <div id="wz-body"></div>
</div>
<div id="wz-toast" class="wz-toast"></div>
<script>window.__WZ_CFG__ = {cfg_json};</script>
{WIZARD_SCRIPT}
</body>
</html>
"""


def _error_shell(message: str) -> str:
    safe = htmlmod.escape(message)
    return f"""<!DOCTYPE html><html><head><meta charset="utf-8">
<style>{THEME_CSS}</style></head><body>
<div class="ibm-masthead"><div class="brand"><span class="dot"></span>IBM Consulting Advantage · AI Reconciliation</div>
<div class="meta">Input required</div></div>
<div class="ibm-shell"><div class="ibm-card"><header>Unable to reconcile</header>
<div class="body"><p>{safe}</p>
<p style="color:var(--ibm-text-2);font-size:12px;margin-top:12px;">Tip: drop both files into the chat using the paperclip icon, then ask the assistant to reconcile them.</p>
</div></div></div></body></html>"""


def _shell_html(
    meta: Dict[str, Any],
    server_fallback: Dict[str, Dict[str, str]],
    inline_payload: Optional[Dict[str, Any]] = None,
) -> str:
    fallback_json = json.dumps(server_fallback)
    payload_json = json.dumps(inline_payload) if inline_payload else "null"
    observer = (
        OBSERVER_SCRIPT
        .replace("__IBM_NAVY__", IBM_NAVY)
        .replace("__IBM_BLUE__", IBM_BLUE)
        .replace("__BUILD__", _BUILD)
    )
    safe_title = htmlmod.escape(meta.get("title", "AI Reconciliation"))
    run_id = htmlmod.escape(meta.get("runId", ""))
    sector = htmlmod.escape(meta.get("sector", ""))
    region = htmlmod.escape(meta.get("region", ""))

    return f"""<!DOCTYPE html>
<html lang="en" data-ibm-build="{_BUILD}">
<head>
<meta charset="utf-8">
<title>{safe_title}</title>
<style>{THEME_CSS}</style>
{CDN_SCRIPTS}
</head>
<body>
<div class="ibm-masthead">
  <div class="brand"><span class="dot"></span>IBM Consulting Advantage · AI Reconciliation</div>
  <div class="meta">{sector} · {region} · run {run_id}</div>
</div>
<div class="ibm-dl-bar">
  <span class="label">Download reconciled file:</span>
  <button class="ibm-btn" id="ibm-dl-xlsx" data-ibm-dl="xlsx" disabled>XLSX</button>
  <button class="ibm-btn" id="ibm-dl-docx" data-ibm-dl="docx" disabled>DOCX</button>
  <button class="ibm-btn" id="ibm-dl-pptx" data-ibm-dl="pptx" disabled>PPTX</button>
  <span class="status" id="ibm-status">Awaiting reconciliation stream…</span>
</div>
<div class="ibm-shell">
  <div id="ibm-render"></div>
  <div id="ibm-loader" class="ibm-loading">
    <div class="dots"><span></span><span></span><span></span></div>
    <div style="margin-top:12px;">Receiving reconciliation output…</div>
  </div>
</div>
<div id="ibm-toast" class="ibm-toast"></div>
<script>
  window.__IBM_SERVER_FALLBACK__ = {fallback_json};
  window.__IBM_RECON_PAYLOAD__ = {payload_json};
</script>
{observer}
</body>
</html>
"""


# ---------------------------------------------------------------------------
# Reconciliation math (deterministic, runs in Python not LLM)
# ---------------------------------------------------------------------------


def _to_float(v: Any) -> Optional[float]:
    if v is None or v == "":
        return None
    try:
        return float(str(v).replace(",", "").strip())
    except (ValueError, TypeError):
        return None


def _norm_str(v: Any) -> str:
    return re.sub(r"\s+", " ", str(v or "").strip()).lower()


def _key_for(rec: Dict[str, Any], key_fields: List[str]) -> str:
    parts = []
    for k in key_fields:
        parts.append(_norm_str(rec.get(k, "")))
    return "||".join(parts)


def _redact_pii(records: List[Dict[str, Any]], sector: str) -> List[Dict[str, Any]]:
    """HIPAA / GDPR-aware redaction for healthcare downloads."""
    if sector.lower() not in ("healthcare", "pharma_clinical", "pharma/clinical", "pharmaceutical"):
        return records
    out = []
    name_keys = {"patientname", "patient", "insured", "policyholder", "subject", "patientfullname"}
    dob_keys = {"patientdob", "dob", "birthdate", "dateofbirth"}
    id_keys = {"ssn", "socialsecurity", "memberid", "member_id", "patientid", "mrn"}
    for r in records:
        new = {}
        for k, v in r.items():
            lk = re.sub(r"[^a-z]", "", k.lower())
            if lk in name_keys and v:
                parts = re.split(r"[\s,]+", str(v).strip())
                masked = " ".join((p[0] + "." if p else "") for p in parts if p)
                new[k] = masked or "***"
            elif lk in dob_keys and v:
                s = str(v)
                new[k] = s[:4] + "-**-**" if len(s) >= 4 else "****"
            elif lk in id_keys and v:
                s = str(v)
                new[k] = (s[:3] + "***" + s[-2:]) if len(s) > 5 else "***"
            else:
                new[k] = v
        out.append(new)
    return out


def _reconcile(
    left: List[Dict[str, Any]],
    right: List[Dict[str, Any]],
    key_fields: List[str],
    amount_fields: List[str],
    tolerance: float,
) -> Dict[str, Any]:
    """Deterministic partition: matched / variance / unmatched_left / unmatched_right."""
    right_index: Dict[str, List[int]] = {}
    for i, r in enumerate(right):
        k = _key_for(r, key_fields)
        right_index.setdefault(k, []).append(i)

    used_right = set()
    matched, variance = [], []
    unmatched_left = []
    for lrec in left:
        k = _key_for(lrec, key_fields)
        cand = right_index.get(k, [])
        picked = None
        for ri in cand:
            if ri not in used_right:
                picked = ri
                break
        if picked is None:
            unmatched_left.append(lrec)
            continue
        used_right.add(picked)
        rrec = right[picked]
        # Check amount fields for variance
        diffs = {}
        for f in amount_fields:
            lv = _to_float(lrec.get(f))
            # right may have a differently-named amount field; check same name first, then common synonyms
            rv = _to_float(rrec.get(f))
            if rv is None:
                for alt in (f.lower(), f.upper(), f.replace("_", "").lower(), "amount", "Amount", "NetAmount", "Consideration", "MV", "MarketValue", "Paid", "PaidAmount", "BilledAmount", "Charged"):
                    if alt in rrec:
                        rv = _to_float(rrec.get(alt))
                        if rv is not None:
                            break
            if lv is None and rv is None:
                continue
            if lv is None or rv is None:
                diffs[f] = {"left": lv, "right": rv, "delta": None}
                continue
            if abs(lv - rv) > tolerance:
                diffs[f] = {"left": lv, "right": rv, "delta": round(rv - lv, 4)}
        if diffs:
            row = {"_key": k}
            row.update(lrec)
            for f, d in diffs.items():
                row[f"{f}_left"] = d["left"]
                row[f"{f}_right"] = d["right"]
                row[f"{f}_delta"] = d["delta"]
            variance.append(row)
        else:
            matched.append(lrec)

    unmatched_right = [right[i] for i in range(len(right)) if i not in used_right]

    total = len(matched) + len(variance) + len(unmatched_left)
    match_rate = round((len(matched) / total * 100), 2) if total else 0.0

    return {
        "matched": matched,
        "variance": variance,
        "unmatched_left": unmatched_left,
        "unmatched_right": unmatched_right,
        "stats": {
            "total_left": len(left),
            "total_right": len(right),
            "matched": len(matched),
            "variance": len(variance),
            "unmatched_left": len(unmatched_left),
            "unmatched_right": len(unmatched_right),
            "match_rate": match_rate,
        },
    }


# ---------------------------------------------------------------------------
# ooXML fallback — always built server-side, shipped as base64 data URIs
# ---------------------------------------------------------------------------


def _build_xlsx_bytes(payload: Dict[str, Any]) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    navy = "FF002D74"
    header_font = Font(bold=True, color="FFFFFFFF", name="Calibri", size=11)
    header_fill = PatternFill("solid", fgColor=navy)
    title_font = Font(bold=True, size=14, color="FF002D74")

    def add_sheet(name: str, rows: List[Dict[str, Any]], title: str) -> None:
        if name == "Summary":
            ws = wb.active
            ws.title = name
        else:
            ws = wb.create_sheet(name)
        ws["A1"] = title
        ws["A1"].font = title_font
        ws["A2"] = "IBM Consulting Advantage — AI Reconciliation"
        if not rows:
            ws["A4"] = "(no rows)"
            return
        keys = list(rows[0].keys())
        for j, k in enumerate(keys, 1):
            c = ws.cell(row=4, column=j, value=k)
            c.font = header_font
            c.fill = header_fill
            c.alignment = Alignment(horizontal="left")
        for i, r in enumerate(rows, 5):
            for j, k in enumerate(keys, 1):
                v = r.get(k, "")
                if isinstance(v, (dict, list)):
                    v = json.dumps(v, default=str)
                ws.cell(row=i, column=j, value=v)
        for j, k in enumerate(keys, 1):
            col = get_column_letter(j)
            ws.column_dimensions[col].width = min(max(len(str(k)), 12), 40)

    meta = payload.get("meta", {})
    stats = payload.get("stats", {})
    summary_rows = [
        {"Field": "Sector", "Value": meta.get("sector", "")},
        {"Field": "Region", "Value": meta.get("region", "")},
        {"Field": "As of", "Value": meta.get("asOf", "")},
        {"Field": "Run ID", "Value": meta.get("runId", "")},
        {"Field": "Regulations", "Value": ", ".join(meta.get("regulations", []))},
    ]
    for k, v in stats.items():
        summary_rows.append({"Field": k, "Value": v})
    add_sheet("Summary", summary_rows, "Reconciliation Summary")
    sections = [
        ("Variance", payload.get("variance", [])),
        ("Missing_in_Right", payload.get("unmatched_left", [])),
        ("Missing_in_Left", payload.get("unmatched_right", [])),
        ("Matched", payload.get("matched", [])),
        ("Commentary", [{"severity": n.get("severity", ""), "regulation": n.get("regulation", ""), "text": n.get("text", "")} for n in payload.get("narrative", [])]),
    ]
    for name, rows in sections:
        add_sheet(name, rows, name.replace("_", " "))
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _build_docx_bytes(payload: Dict[str, Any]) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()
    meta = payload.get("meta", {})
    stats = payload.get("stats", {})

    # Title
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run("IBM Consulting Advantage")
    run.bold = True
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x00, 0x2D, 0x74)

    h = doc.add_heading("AI Reconciliation Report", level=0)
    for r in h.runs:
        r.font.color.rgb = RGBColor(0x00, 0x2D, 0x74)

    doc.add_paragraph(f"Sector: {meta.get('sector','')}   Region: {meta.get('region','')}   As of: {meta.get('asOf','')}")
    if meta.get("regulations"):
        doc.add_paragraph("Regulations: " + ", ".join(meta["regulations"]))
    if meta.get("redacted"):
        doc.add_paragraph("Note: patient/subject PII has been redacted in accordance with HIPAA / GDPR.")

    doc.add_heading("Summary", level=1)
    tbl = doc.add_table(rows=1, cols=2)
    tbl.style = "Light Grid Accent 1"
    hdr = tbl.rows[0].cells
    hdr[0].text = "KPI"
    hdr[1].text = "Value"
    for k, v in stats.items():
        row = tbl.add_row().cells
        row[0].text = str(k)
        row[1].text = str(v)

    def add_table(title: str, rows: List[Dict[str, Any]]) -> None:
        doc.add_heading(title + f" ({len(rows)})", level=1)
        if not rows:
            doc.add_paragraph("(none)")
            return
        keys = list(rows[0].keys())[:8]
        t = doc.add_table(rows=1, cols=len(keys))
        t.style = "Light Grid Accent 1"
        for j, k in enumerate(keys):
            t.rows[0].cells[j].text = k
        for r in rows[:200]:
            cells = t.add_row().cells
            for j, k in enumerate(keys):
                v = r.get(k, "")
                if isinstance(v, (dict, list)):
                    v = json.dumps(v, default=str)
                cells[j].text = str(v)[:120]

    add_table("Variances", payload.get("variance", []))
    add_table("Missing in Right", payload.get("unmatched_left", []))
    add_table("Missing in Left", payload.get("unmatched_right", []))
    if payload.get("matched"):
        add_table("Matched (sample)", payload.get("matched", [])[:50])

    if payload.get("narrative"):
        doc.add_heading("Variance commentary", level=1)
        for n in payload["narrative"]:
            para = doc.add_paragraph()
            if n.get("regulation"):
                r1 = para.add_run(f"[{n['regulation']}] ")
                r1.bold = True
                r1.font.color.rgb = RGBColor(0x00, 0x43, 0xCE)
            para.add_run(n.get("text", ""))

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _build_pptx_bytes(payload: Dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    meta = payload.get("meta", {})
    stats = payload.get("stats", {})
    NAVY = RGBColor(0x00, 0x2D, 0x74)
    BLUE = RGBColor(0x00, 0x43, 0xCE)
    SURF = RGBColor(0xED, 0xF5, 0xFF)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    blank = pres.slide_layouts[6]

    def title_bar(slide, text):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, pres.slide_width, Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.05), Inches(12), Inches(0.5))
        tf = tb.text_frame
        tf.text = text
        tf.paragraphs[0].runs[0].font.size = Pt(20)
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.color.rgb = WHITE

    # Slide 1: cover
    s = pres.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, pres.slide_width, pres.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.6))
    tb.text_frame.text = "IBM Consulting Advantage"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    tb.text_frame.paragraphs[0].runs[0].font.bold = True
    tb.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(1.2))
    tb.text_frame.text = "AI Reconciliation Report"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(44)
    tb.text_frame.paragraphs[0].runs[0].font.bold = True
    tb.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(0.5))
    tb.text_frame.text = f"{meta.get('sector','')} · {meta.get('region','')} · as of {meta.get('asOf','')}"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    tb.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xD0, 0xE2, 0xFF)

    # Slide 2: KPIs
    s = pres.slides.add_slide(blank)
    title_bar(s, "Reconciliation KPIs")
    items = list(stats.items())
    per = 4
    cw, ch = Inches(2.9), Inches(1.1)
    for i, (k, v) in enumerate(items):
        col, row = i % per, i // per
        x = Inches(0.5 + col * 3.0)
        y = Inches(1.0 + row * 1.3)
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, ch)
        box.fill.solid()
        box.fill.fore_color.rgb = SURF
        box.line.color.rgb = RGBColor(0xD0, 0xE2, 0xFF)
        tb = s.shapes.add_textbox(x, y, cw, Inches(0.6))
        tb.text_frame.text = str(v)
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(22)
        tb.text_frame.paragraphs[0].runs[0].font.bold = True
        tb.text_frame.paragraphs[0].runs[0].font.color.rgb = NAVY
        tb2 = s.shapes.add_textbox(x, Emu(y + Inches(0.55)), cw, Inches(0.5))
        tb2.text_frame.text = str(k)
        tb2.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
        tb2.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x52, 0x52, 0x52)

    def table_slide(title: str, rows: List[Dict[str, Any]]) -> None:
        if not rows:
            return
        s = pres.slides.add_slide(blank)
        title_bar(s, f"{title} ({len(rows)})")
        keys = list(rows[0].keys())[:6]
        n_rows = min(len(rows), 14) + 1
        tbl_shape = s.shapes.add_table(n_rows, len(keys), Inches(0.3), Inches(0.8), Inches(12.7), Inches(6.2))
        tbl = tbl_shape.table
        for j, k in enumerate(keys):
            cell = tbl.cell(0, j)
            cell.text = str(k)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = WHITE
                    r.font.bold = True
                    r.font.size = Pt(10)
        for i, r in enumerate(rows[:14], 1):
            for j, k in enumerate(keys):
                v = r.get(k, "")
                if isinstance(v, (dict, list)):
                    v = json.dumps(v, default=str)
                cell = tbl.cell(i, j)
                cell.text = str(v)[:60]
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(9)

    table_slide("Variances", payload.get("variance", []))
    table_slide("Missing in Right", payload.get("unmatched_left", []))
    table_slide("Missing in Left", payload.get("unmatched_right", []))

    # Narrative
    if payload.get("narrative"):
        s = pres.slides.add_slide(blank)
        title_bar(s, "Variance Commentary")
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12.3), Inches(6.2))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for n in payload["narrative"][:10]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = 0
            tag = f"[{n.get('regulation','')}] " if n.get("regulation") else ""
            run1 = p.add_run()
            run1.text = tag
            run1.font.bold = True
            run1.font.color.rgb = BLUE
            run1.font.size = Pt(14)
            run2 = p.add_run()
            run2.text = n.get("text", "")
            run2.font.size = Pt(14)
            run2.font.color.rgb = RGBColor(0x16, 0x16, 0x16)

    buf = io.BytesIO()
    pres.save(buf)
    return buf.getvalue()


def _b64_data_uri(mime: str, data: bytes) -> str:
    return f"data:{mime};base64,{base64.b64encode(data).decode('ascii')}"


def _build_server_fallback(payload: Dict[str, Any], base_name: str) -> Dict[str, Dict[str, str]]:
    out: Dict[str, Dict[str, str]] = {}
    try:
        xlsx = _build_xlsx_bytes(payload)
        out["xlsx"] = {
            "filename": f"{base_name}.xlsx",
            "dataUri": _b64_data_uri(
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", xlsx
            ),
        }
    except Exception as e:
        out["xlsx_error"] = {"filename": "", "dataUri": f"data:text/plain;base64,{base64.b64encode(str(e).encode()).decode()}"}
    try:
        docx = _build_docx_bytes(payload)
        out["docx"] = {
            "filename": f"{base_name}.docx",
            "dataUri": _b64_data_uri(
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document", docx
            ),
        }
    except Exception as e:
        out["docx_error"] = {"filename": "", "dataUri": f"data:text/plain;base64,{base64.b64encode(str(e).encode()).decode()}"}
    try:
        pptx = _build_pptx_bytes(payload)
        out["pptx"] = {
            "filename": f"{base_name}.pptx",
            "dataUri": _b64_data_uri(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation", pptx
            ),
        }
    except Exception as e:
        out["pptx_error"] = {"filename": "", "dataUri": f"data:text/plain;base64,{base64.b64encode(str(e).encode()).decode()}"}
    return out


# ---------------------------------------------------------------------------
# Tool class
# ---------------------------------------------------------------------------


class Tools:
    """IBM Consulting Advantage — AI Reconciliation.

    The LLM parses uploaded regulatory documents from Open WebUI's attachment
    context, normalises them into two structured record lists (left / right),
    decides the key fields and amount fields per sector, and calls
    'reconcile' to get a deterministic partition plus an inline IBM-themed
    report with CDN-first / ooXML-fallback downloads (XLSX / DOCX / PPTX).
    """

    class Valves(BaseModel):
        default_tolerance: float = Field(
            default=0.01,
            description="Default numeric tolerance (absolute) for amount comparisons when the model does not specify one.",
        )
        max_preview_rows: int = Field(
            default=200,
            description="Maximum rows per section included in the inline preview and downloadable outputs.",
        )

    def __init__(self) -> None:
        self.valves = self.Valves()

    async def reconcile(
        self,
        sector: Optional[str] = None,
        region: Optional[str] = None,
        as_of: Optional[str] = None,
        output_name: Optional[str] = None,
        tolerance: Optional[float] = None,
        __event_call__=None,
        __event_emitter__=None,
    ):
        """Launch the AI Reconciliation wizard in an inline iframe.

        Call this tool whenever the user expresses a reconciliation intent —
        uploading files, comparing two datasets, month-end close, etc.
        The iframe renders a stepwise wizard that lets the user upload two
        datasets (CSV / TSV / XLSX / XLS / DOCX / PPTX / PDF / JSON),
        configure matching rules, run reconciliation in the browser, and
        download XLSX / CSV / PDF / DOCX outputs with a full audit trail
        (run ID, SHA-256 of inputs, operator name, regulation tags).

        All file parsing, matching, and output generation runs client-side
        for speed. No records round-trip to Python. CSV-only fallback
        kicks in automatically when CDN scripts are blocked.

        :param sector: Optional sector hint to pre-select (Banking, Investment Banking, Insurance, Healthcare, Asset Management, Pharma Clinical, Energy, Telecommunications, Retail, Manufacturing, Public Sector, Technology, Transportation, Other).
        :param region: Optional region hint to pre-select (USA, European Union, United Kingdom, Global).
        :param as_of: Optional reporting date YYYY-MM-DD.
        :param output_name: Optional base filename for downloads.
        :param tolerance: Optional absolute tolerance for numeric equality.
        :return: Inline HTMLResponse rendering the wizard.
        """
        log.info("[recon] launching wizard: sector=%r region=%r", sector, region)
        if __event_emitter__:
            await __event_emitter__({"type": "status", "data": {"description": "Loading reconciliation wizard...", "done": False}})
        try:
            html = _wizard_shell_html()
            if __event_emitter__:
                await __event_emitter__({"type": "status", "data": {"description": "Wizard ready.", "done": True}})
            return HTMLResponse(content=html, headers={"Content-Disposition": "inline"})
        except Exception as exc:
            tb = traceback.format_exc()
            log.error("[recon] wizard render failed: %s\n%s", exc, tb)
            if __event_emitter__:
                await __event_emitter__({"type": "status", "data": {"description": f"Error: {type(exc).__name__}", "done": True}})
            return HTMLResponse(
                content=_error_shell(
                    f"Wizard render failed with: {type(exc).__name__}: {exc}."
                ),
                headers={"Content-Disposition": "inline"},
            )


async def _do_reconcile(
    valves,
    left,
    right,
    sector,
    region,
    key_fields,
    amount_fields,
    tolerance,
    regulations,
    as_of,
    output_name,
    narrative,
    event_emitter,
):
    sector = sector or "Banking"
    region = region or "USA"
    # Normalise legacy region spellings
    region_map = {"Europe": "European Union", "EU": "European Union", "UK": "United Kingdom"}
    region = region_map.get(region, region)
    defaults = _sector_defaults(sector)
    kf = key_fields or defaults["key_fields"]
    af = amount_fields or defaults["amount_fields"]
    tol = float(tolerance) if tolerance is not None else valves.default_tolerance
    reg_table = defaults.get("regulations", {}) if isinstance(defaults, dict) else {}
    regs = regulations or reg_table.get(region) or reg_table.get("USA") or []

    if event_emitter:
        await event_emitter({"type": "status", "data": {"description": f"Matching {len(left)} vs {len(right)} records on {', '.join(kf)}…", "done": False}})

    result = _reconcile(left, right, kf, af, tol)

    redacted = False
    if sector.lower() in ("healthcare", "pharma clinical", "pharma_clinical", "pharmaceutical"):
        for bucket in ("matched", "variance", "unmatched_left", "unmatched_right"):
            result[bucket] = _redact_pii(result[bucket], sector)
        redacted = True

    cap = valves.max_preview_rows
    for bucket in ("matched", "variance", "unmatched_left", "unmatched_right"):
        if len(result[bucket]) > cap:
            result[bucket] = result[bucket][:cap]

    run_id = uuid.uuid4().hex[:8]
    meta = {
        "title": f"AI Reconciliation — {sector} ({region})",
        "sector": sector,
        "region": region,
        "regulations": regs,
        "asOf": as_of or time.strftime("%Y-%m-%d"),
        "runId": run_id,
        "outputName": output_name or f"reconciliation_{sector.lower().replace(' ', '_')}",
        "redacted": redacted,
    }
    payload = {
        "meta": meta,
        "stats": result["stats"],
        "matched": result["matched"],
        "variance": result["variance"],
        "unmatched_left": result["unmatched_left"],
        "unmatched_right": result["unmatched_right"],
        "narrative": narrative or [],
    }

    if event_emitter:
        await event_emitter({"type": "status", "data": {"description": "Building ooXML server fallback (XLSX / DOCX / PPTX)…", "done": False}})

    base_name = f"{meta['outputName']}_{meta['asOf'].replace('-', '')}"
    server_fallback = _build_server_fallback(payload, base_name)

    html = _shell_html(meta, server_fallback, inline_payload=payload)

    if event_emitter:
        await event_emitter({
            "type": "status",
            "data": {
                "description": (
                    f"Reconciliation complete — {result['stats']['match_rate']}% match rate · "
                    f"{result['stats']['variance']} variance · "
                    f"{result['stats']['unmatched_left']+result['stats']['unmatched_right']} unmatched."
                ),
                "done": True,
            },
        })
    return HTMLResponse(content=html, headers={"Content-Disposition": "inline"})


# ---------------------------------------------------------------------------
# Sector defaults — key fields, amount fields, governing regulations
# ---------------------------------------------------------------------------


def _sector_defaults(sector: str) -> Dict[str, Any]:
    s = (sector or "").lower().strip()
    if s == "banking":
        return {
            "key_fields": ["TransactionID"],
            "amount_fields": ["Amount"],
            "regulations": {
                "USA": ["FFIEC 031", "Reg W", "BCBS 239", "SOX 404"],
                "European Union": ["CRR/CRD IV", "FINREP", "BCBS 239", "DORA"],
                "United Kingdom": ["PRA SS1/23", "BCBS 239", "UK GAAP"],
                "Global": ["BCBS 239", "Basel III"],
            },
        }
    if s == "investment banking":
        return {
            "key_fields": ["TradeID"],
            "amount_fields": ["NetAmount", "Quantity", "Price"],
            "regulations": {
                "USA": ["SEC 10-Q", "FINRA TRACE", "Dodd-Frank", "CFTC Part 45"],
                "European Union": ["MiFID II RTS 22", "EMIR REFIT", "CSDR"],
                "United Kingdom": ["FCA MAR", "UK EMIR"],
                "Global": ["IOSCO", "CPMI-IOSCO PFMI"],
            },
        }
    if s == "insurance":
        return {
            "key_fields": ["PolicyNumber"],
            "amount_fields": ["GrossPremium", "NetPremium", "ClaimAmount", "Paid"],
            "regulations": {
                "USA": ["NAIC SAP", "ORSA", "SOX 404"],
                "European Union": ["Solvency II", "IFRS 17", "EIOPA"],
                "United Kingdom": ["PRA SS4/18", "FCA ICOBS"],
                "Global": ["IFRS 17", "IAIS ICP"],
            },
        }
    if s == "healthcare":
        return {
            "key_fields": ["ClaimID"],
            "amount_fields": ["BilledAmount", "Paid", "Allowed", "Charged"],
            "regulations": {
                "USA": ["HIPAA 837/835", "CMS-1500", "HITECH"],
                "European Union": ["GDPR", "EHDS", "MDR"],
                "United Kingdom": ["UK GDPR", "NHS DSPT"],
                "Global": ["HL7 FHIR", "ICD-10"],
            },
        }
    if s == "asset management":
        return {
            "key_fields": ["AccountID", "ISIN"],
            "amount_fields": ["Quantity", "MarketValue"],
            "regulations": {
                "USA": ["SEC 13F", "Investment Company Act", "Form PF"],
                "European Union": ["UCITS", "AIFMD", "SFDR"],
                "United Kingdom": ["FCA COLL", "UK UCITS"],
                "Global": ["IOSCO"],
            },
        }
    if s in ("pharma clinical", "pharma_clinical", "pharmaceutical", "pharma/clinical"):
        return {
            "key_fields": ["SubjectID", "VisitDate"],
            "amount_fields": ["SystolicBP", "DiastolicBP", "HeartRate"],
            "regulations": {
                "USA": ["FDA 21 CFR Part 11", "ICH GCP E6(R3)"],
                "European Union": ["EMA EudraCT", "GDPR", "CTR"],
                "United Kingdom": ["MHRA GCP", "UK GDPR"],
                "Global": ["ICH GCP", "CDISC"],
            },
        }
    if s in ("energy", "energy & utilities", "utilities"):
        return {
            "key_fields": ["MeterID", "ReadingDate"],
            "amount_fields": ["UsagekWh", "BilledAmount"],
            "regulations": {
                "USA": ["FERC Order 2222", "NERC CIP", "SOX 404"],
                "European Union": ["REMIT", "ESRS E1", "EU ETS"],
                "United Kingdom": ["Ofgem", "UK ETS"],
                "Global": ["ISO 50001"],
            },
        }
    if s in ("telecommunications", "telco", "telecom"):
        return {
            "key_fields": ["SubscriberID", "BillingCycle"],
            "amount_fields": ["ChargeAmount", "Usage"],
            "regulations": {
                "USA": ["FCC CPNI", "SOX 404"],
                "European Union": ["EECC", "GDPR", "NIS2"],
                "United Kingdom": ["Ofcom GC"],
                "Global": ["TM Forum"],
            },
        }
    if s in ("retail", "retail & consumer", "consumer"):
        return {
            "key_fields": ["OrderID"],
            "amount_fields": ["NetAmount", "TaxAmount", "TotalAmount"],
            "regulations": {
                "USA": ["SOX 404", "PCI DSS", "GAAP"],
                "European Union": ["CSRD/ESRS", "PSD2", "GDPR"],
                "United Kingdom": ["UK GAAP", "PSR"],
                "Global": ["IFRS 15", "PCI DSS"],
            },
        }
    if s in ("manufacturing", "manufacturing & industrial", "industrial"):
        return {
            "key_fields": ["POID", "MaterialCode"],
            "amount_fields": ["Quantity", "UnitPrice", "LineAmount"],
            "regulations": {
                "USA": ["SOX 404", "FDA 21 CFR 820", "GAAP"],
                "European Union": ["CSRD/ESRS", "CE Marking", "REACH"],
                "United Kingdom": ["UKCA", "UK GAAP"],
                "Global": ["ISO 9001", "IFRS 15"],
            },
        }
    if s in ("public sector", "government", "public"):
        return {
            "key_fields": ["VoucherID"],
            "amount_fields": ["ObligatedAmount", "DisbursedAmount"],
            "regulations": {
                "USA": ["FAR/DFARS", "GAO Green Book", "NIST 800-53"],
                "European Union": ["EU Financial Regulation", "ESIF"],
                "United Kingdom": ["Managing Public Money", "FReM"],
                "Global": ["IPSAS"],
            },
        }
    if s in ("technology", "software", "technology / software"):
        return {
            "key_fields": ["InvoiceID"],
            "amount_fields": ["Subtotal", "Tax", "Total"],
            "regulations": {
                "USA": ["SOX 404", "GAAP ASC 606"],
                "European Union": ["CSRD/ESRS", "DORA", "GDPR"],
                "United Kingdom": ["UK GAAP", "UK DPA"],
                "Global": ["IFRS 15", "ISO 27001"],
            },
        }
    if s in ("transportation", "logistics", "transportation & logistics"):
        return {
            "key_fields": ["ShipmentID"],
            "amount_fields": ["FreightAmount", "DutyAmount", "TotalAmount"],
            "regulations": {
                "USA": ["DOT FMCSA", "CBP 19 CFR", "SOX 404"],
                "European Union": ["UCC", "CSRD/ESRS"],
                "United Kingdom": ["HMRC CDS"],
                "Global": ["IATA", "IMO"],
            },
        }
    return {
        "key_fields": ["ID"],
        "amount_fields": ["Amount"],
        "regulations": {
            "USA": ["SOX 404", "GAAP"],
            "European Union": ["IFRS", "GDPR"],
            "United Kingdom": ["UK GAAP", "UK GDPR"],
            "Global": ["IFRS", "ISO 27001"],
        },
    }
